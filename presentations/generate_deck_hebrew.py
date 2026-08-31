#!/usr/bin/env python3
"""
יצירת מצגת פתרון ארגונית מקיפה בעברית (.pptx)
עבור פלטפורמת ניהול המסמכים ב-AWS (AWS Document Management Platform).
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -----------------------------------------------------------------------------
# ערכת צבעים וטיפוגרפיה
# -----------------------------------------------------------------------------
BG_COLOR = RGBColor(11, 17, 32)       # #0B1120 Deep Midnight Navy
CARD_BG = RGBColor(30, 41, 59)        # #1E293B Slate Dark Card
CARD_BORDER = RGBColor(51, 65, 85)    # #334155 Slate Border
CARD_BG_LIGHT = RGBColor(24, 32, 47)  # #18202F

TEXT_WHITE = RGBColor(248, 250, 252)  # #F8FAFC
TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8
TEXT_DIM = RGBColor(100, 116, 139)    # #64748B

AWS_ORANGE = RGBColor(255, 153, 0)    # #FF9900 AWS Orange
ACCENT_BLUE = RGBColor(56, 189, 248)  # #38BDF8 Sky Blue
ACCENT_INDIGO = RGBColor(129, 140, 248) # #818CF8 Indigo
ACCENT_GREEN = RGBColor(52, 211, 153) # #34D399 Emerald Green
ACCENT_PURPLE = RGBColor(192, 132, 252) # #C084FC Purple
ACCENT_ROSE = RGBColor(244, 63, 94)   # #F43F5E Rose Red
ACCENT_AMBER = RGBColor(251, 191, 36) # #FBBF24 Amber

FONT_FAMILY = "Segoe UI"
FONT_FAMILY_TITLE = "Segoe UI"

def set_slide_background(slide, color=BG_COLOR):
    """צביעת רקע השקף בצבע אחיד."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, subtitle_text, category_tag="ארכיטקטורת AWS Serverless"):
    """הוספת כותרת ראשית ומודרנית לשקף."""
    # תגית קטגוריה
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = 0
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = category_tag.upper()
    p_tag.font.name = FONT_FAMILY
    p_tag.font.size = Pt(9.5)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_BLUE

    # כותרת שקף
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(11.7), Inches(0.55))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_FAMILY_TITLE
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

    # כותרת משנה
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.22), Inches(11.7), Inches(0.35))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.name = FONT_FAMILY
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = TEXT_MUTED

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    """יצירת כרטיס מעוצב."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape

def add_footer(slide, current_slide, total_slides):
    """הוספת קו תחתון ומספור שקפים."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CARD_BORDER
    line.line.fill.background()

    box_l = slide.shapes.add_textbox(Inches(0.8), Inches(7.12), Inches(8.5), Inches(0.25))
    tf_l = box_l.text_frame
    tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
    p_l = tf_l.paragraphs[0]
    p_l.text = "פלטפורמת ניהול מסמכים ארגונית ב-AWS | ארכיטקטורת Serverless ו-WORM"
    p_l.font.name = FONT_FAMILY
    p_l.font.size = Pt(8.5)
    p_l.font.color.rgb = TEXT_DIM

    box_r = slide.shapes.add_textbox(Inches(10.5), Inches(7.12), Inches(2.033), Inches(0.25))
    tf_r = box_r.text_frame
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
    p_r = tf_r.paragraphs[0]
    p_r.text = f"שקף {current_slide} מתוך {total_slides}"
    p_r.alignment = PP_ALIGN.RIGHT
    p_r.font.name = FONT_FAMILY
    p_r.font.size = Pt(8.5)
    p_r.font.bold = True
    p_r.font.color.rgb = ACCENT_BLUE


def build_hebrew_presentation(output_path="AWS_Document_Management_Platform_Architecture_Hebrew.pptx"):
    script_dir = os.path.dirname(os.path.abspath(__file__))
    if not os.path.isabs(output_path):
        output_path = os.path.join(script_dir, output_path)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    TOTAL_SLIDES = 23

    # =========================================================================
    # SLIDE 1: Title Slide (Executive Dark)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, BG_COLOR)

    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(1.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AWS_ORANGE
    bar.line.fill.background()

    tag_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.4))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ארכיטקטורת פתרון ארגוני ב-AWS"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(11.5), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "פלטפורמת ניהול מסמכים ארגונית בענן (Cloud-Native)"
    p.font.name = FONT_FAMILY_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    sub_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.2))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ארכיטקטורת Serverless מבוזרת, מודל סמכות משולש (Tri-Partite Source of Truth), אימות WORM ועמידה בתקני רגולציה מחמירים"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED

    # 4 כרטיסי היילייט
    card_y = Inches(4.3)
    card_w = Inches(2.78)
    card_h = Inches(2.3)
    card_gap = Inches(0.2)
    left_start = Inches(0.8)

    highlights = [
        ("100% Serverless", "חיסכון מלא בעלויות סרק: S3, DynamoDB On-Demand, Lambda, SQS ו-OpenSearch Serverless ללא שרתים מנוהלים.", ACCENT_BLUE),
        ("מודל סמכות משולש", "הפרדה חדה בין תוכן בינארי (S3), מטא-דאטה מובנה (S3 Annotations), ומנוע בקרה (DynamoDB OCC).", ACCENT_GREEN),
        ("זמני תגובה מהירים", "קליטת קבצים ב-~150ms, שאילתות DynamoDB ב-sub-10ms, והורדות ישירות מ-S3 ב-sub-50ms.", ACCENT_PURPLE),
        ("אבטחה וביקורת WORM", "הצפנת KMS ייעודית, בקרת הרשאות RBAC ב-Cognito, יומן ביקורת בלתי-הפיך ותיוג משאבים מלא.", ACCENT_ROSE)
    ]

    for i, (title, desc, color) in enumerate(highlights):
        cx = left_start + i * (card_w + card_gap)
        add_card(slide1, cx, card_y, card_w, card_h, CARD_BG, CARD_BORDER)

        t_box = slide1.shapes.add_textbox(cx + Inches(0.2), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.35))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color

        d_box = slide1.shapes.add_textbox(cx + Inches(0.2), card_y + Inches(0.6), card_w - Inches(0.4), Inches(1.5))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED

    add_footer(slide1, 1, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 2: Executive Summary & Enterprise Architecture Blueprint
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "תמצית מנהלים וארכיטקטורת על של הפתרון",
               "מודרניזציה של ניהול תוכן ארגוני באמצעות שירותי Serverless מבוזרים וחסרי שרתים בענן",
               "ערך אסטרטגי והקשר עסקי")

    add_card(slide2, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    left_box = slide2.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf = left_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "עקרונות ליבה בארכיטקטורה הארגונית"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ROSE

    items_left = [
        ("אלסטיות מבוססת צריכה בפועל:", "ארכיטקטורת Serverless מלאה מבטיחה אפס עלויות סרק בזמני מנוחה והתרחבות אוטומטית בעומסי שיא."),
        ("מודל סמכות מבוזר ומופרד:", "הפרדה בין אחסון קבצים, אינדוקס מטא-דאטה ומנוע חיפוש טקסטואלי מונעת נעילות מערכתיות וכשלים שרשרתיים."),
        ("S3 Object Annotations מובנה:", "ביטול מגבלת 2 KB של כותרות S3 ללא שימוש בקבצי Sidecar, תוך תמיכה מלאה בהתפתחות סכמות."),
        ("עמידה ברגולציית WORM ושקיפות ביקורת:", "אכיפת גרסאות קבצים בלתי-הפיכות, בקרת מקביליות אופטימית (OCC), ויומן ביקורת רציף ב-S3.")
    ]
    for title, desc in items_left:
        p = tf.add_paragraph()
        p.text = f"• {title} "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(8)

    add_card(slide2, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    right_box = slide2.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf = right_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "יכולות מרכזיות ותוצאות עסקיות"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    items_right = [
        ("עלות בעלות כוללת (TCO) מינימלית:", "מודל תשלום לפי שימוש (Lambda, S3, DynamoDB, OpenSearch Serverless) ללא צורך בהקצאת שרתים יקרה מראש."),
        ("שרידות של 11 תשיעיות (11 9s):", "אחסון קבצים ב-Amazon S3 עם חסימת מחיקה מכוונת באמצעות מדיניות s3:DeleteObjectVersion DENY קשיחה."),
        ("זמני תגובה נמוכים במיוחד:", "העלאת קבצים רגילים ב-~150ms; קישורי הורדה ישירים ב-<50ms; בקרת DynamoDB טרנזקציונית ב-sub-10ms."),
        ("בידוד מוחלט של מרחב הכשל (Blast Radius):", "תקלה זמנית במנוע החיפוש OpenSearch אינה מונעת העלאה, הורדה או עריכת מטא-דאטה של מסמכים.")
    ]
    for title, desc in items_right:
        p = tf.add_paragraph()
        p.text = f"✓ {title} "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(8)

    add_footer(slide2, 2, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 3: End-to-End Solution Architecture Diagram
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "תרשים ארכיטקטורת הפתרון המלאה (End-to-End)",
               "טופולוגיית Serverless מלאה ב-AWS: בקרת כניסה, שכבת מחשוב מבוזרת ושכבות אחסון מופרדות",
               "טופולוגיית המערכת")

    arch_img_path = os.path.join(script_dir, "diagrams/aws_document_management_architecture.png")
    if os.path.exists(arch_img_path):
        slide3.shapes.add_picture(arch_img_path, Inches(0.8), Inches(1.65), Inches(8.2), Inches(5.15))

    add_card(slide3, Inches(9.2), Inches(1.65), Inches(3.333), Inches(5.15), CARD_BG, CARD_BORDER)
    side_box = slide3.shapes.add_textbox(Inches(9.4), Inches(1.85), Inches(2.933), Inches(4.75))
    tf = side_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "עיקרי הארכיטקטורה"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    p_items = [
        ("שער ה-API (API Gateway):", "ניתוב REST מאובטח, אימות JWT באמצעות Cognito, והעברת תוכן בינארי ישיר (*/*).", ACCENT_BLUE),
        ("מיקרו-שירותי Lambda מבודדים:", "פונקציות פיקוד, שאילתה וחיפוש הפועלות תחת הרשאות IAM מינימליות (Least Privilege).", ACCENT_GREEN),
        ("S3 Object Annotations מובנה:", "מטא-דאטה מוסמך בפורמט JSON הצמוד ישירות לגרסת אובייקט ה-S3 (document-metadata).", ACCENT_AMBER),
        ("טבלת בקרת מצב (DynamoDB):", "טבלה מרכזית המנהלת מצבי מסמך (DOC#{id}), היסטוריית גרסאות וסשנים בריצת sub-10ms.", ACCENT_PURPLE),
        ("צינור עיבוד אסינכרוני (CDC):", "זרמי DynamoDB מעבירים שינויים ליומן ביקורת ב-S3 ולתור SQS לאינדוקס ב-OpenSearch.", ACCENT_BLUE)
    ]
    for title, desc, color in p_items:
        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_before = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(9)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(1)

    add_footer(slide3, 3, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 4: The 5-Tier Authority & Consistency Model
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "מודל הסמכות והעקביות ב-5 שכבות (5-Tier Authority Model)",
               "מניעת נעילות טרנזקציוניות מבוזרות באמצעות הגדרת גבולות סמכות ברורים לכל שירות",
               "משילות נתונים וסמכות")

    tiers = [
        ("1. סמכות התוכן (Content Authority)", "Amazon S3 (Primary Bucket)", "אחסון קבצים בינאריים בלתי-הפיך (PDF, TIFF, Office). ניהול גרסאות עם מדיניות נעילה s3:DeleteObjectVersion DENY לעמידה מלאה בתקני WORM. חתימות הצפנה SHA-256.", ACCENT_BLUE),
        ("2. סמכות המטא-דאטה (Metadata Authority)", "Amazon S3 Annotations", "מטא-דאטה מלא ומובנה בפורמט JSON המאוחסן כ-Annotation רשמי (document-metadata) הצמוד ישירות לגרסת ה-S3 ומאומת עפ\"י סכמת bank.document-metadata/1.", ACCENT_GREEN),
        ("3. סמכות הבקרה (Control Plane)", "Amazon DynamoDB", "טבלת בקרת מצב המנהלת מצביעי מסמכים פעילים (DOC#{id}), היסטוריית גרסאות (VER#{pad}), נעילות אופטימיות (OCC) וסשנים בריצת sub-10ms.", AWS_ORANGE),
        ("4. היטל חיפוש נגזר (Search Read Projection)", "OpenSearch Serverless", "מודל קריאה CQRS אסינכרוני המעודכן מזרמי DynamoDB דרך תור SQS. כשל בחיפוש אינו משפיע על העלאה, עריכה או הורדה של מסמכים.", ACCENT_PURPLE),
        ("5. שקיפות ביקורת (Non-Repudiation Audit)", "Amazon S3 Audit Bucket", "יומן שינויים בלתי-הפיך המועבר אוטומטית לדלי ביקורת ייעודי ומוצפן ב-KMS בחלוקה לפי תאריכים לצרכי רגולציה וחקירה משפטית.", ACCENT_ROSE)
    ]

    tier_h = Inches(0.95)
    tier_gap = Inches(0.12)
    start_y = Inches(1.65)

    for i, (title, service, desc, color) in enumerate(tiers):
        ty = start_y + i * (tier_h + tier_gap)
        add_card(slide4, Inches(0.8), ty, Inches(11.733), tier_h, CARD_BG, CARD_BORDER)

        tag_b = slide4.shapes.add_textbox(Inches(1.05), ty + Inches(0.15), Inches(3.2), Inches(0.65))
        tf = tag_b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color

        p_srv = tf.add_paragraph()
        p_srv.text = service
        p_srv.font.name = FONT_FAMILY
        p_srv.font.size = Pt(9.5)
        p_srv.font.color.rgb = TEXT_MUTED

        desc_b = slide4.shapes.add_textbox(Inches(4.4), ty + Inches(0.15), Inches(7.9), Inches(0.65))
        tf_d = desc_b.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_WHITE

    add_footer(slide4, 4, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 5: Interplay: S3 Annotations vs. DynamoDB Control Plane
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "חלוקת תפקידים: S3 Annotations מול טבלת הבקרה ב-DynamoDB",
               "מדוע נדרשים שני השירותים: פתרון מגבלות אחסון אובייקטים באמצעות מהירות של מסד נתונים טרנזקציוני",
               "ארכיטקטורת נתונים מעמיקה")

    rows = [
        ("מימד ארכיטקטוני", "Amazon S3 & S3 Annotations", "טבלת הבקרה ב-DynamoDB", "מדוע DynamoDB חיוני?"),
        ("סמכות מטא-דאטה מלאה", "סמכות עליונה (מכיל את ה-JSON המלא לפי סכמה)", "שומר מצביע קל משקל, מונה עדכונים ו-ETag", "S3 Annotations מצמיד את המטא-דאטה ישירות לגרסת הקובץ ומאפשר שאילתות SQL באגם נתונים."),
        ("בקרת מקביליות (OCC)", "אין תמיכה בהתניות מורכבות על שדות JSON", "סמכות ראשית: אכיפה אטומית ConditionExpression", "DynamoDB מוודא אטומית current_metadata_revision = :exp למניעת דריסת עדכונים מקבילים."),
        ("איתור מצביע פעיל O(1)", "דורש סריקה איטית של כלל גרסאות האובייקט ב-Bucket", "סמכות ראשית: שליפת מפתח ישירה בפחות מ-10ms", "DynamoDB ממפה DOC#{id} ישירות ל-s3_version_id העדכני ללא קריאות ListObjects איטיות."),
        ("היסטוריית גרסאות מסודרת", "מזהי גרסאות ב-S3 אינם עוקבים ואינם קריאים", "סמכות ראשית: שליפה מהירה לפי begins_with(VER#)", "מאפשר שליפת היסטוריית גרסאות ממוינת (v1, v2, v3) בשאילתה אולטרה-מהירה אחת."),
        ("טרנזקציות רב-רכיביות (ACID)", "אין מנגנון טרנזקציוני מובנה רב-אובייקטי ב-S3", "סמכות ראשית: TransactWriteItems אטומי", "מעדכן אטומית את מצביע המסמך, רושם גרסה חדשה ונועל מזהה Idempotency בקריאה אחת."),
        ("קליטת שינויים רציפה (CDC)", "אירועי S3 אינם כוללים תמונת מצב לפני/אחרי", "סמכות ראשית: DynamoDB Streams מובנה ומסודר", "פולט אירועים מסודרים המזינים אסינכרונית את יומן הביקורת ב-S3 ואת מנוע החיפוש OpenSearch.")
    ]

    table_shape = slide5.shapes.add_table(len(rows), 4, Inches(0.8), Inches(1.65), Inches(11.733), Inches(5.1))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.2)
    table.columns[3].width = Inches(3.333)

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(22, 30, 46)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.12)
            tf.margin_top = tf.margin_bottom = Inches(0.08)

            p = tf.paragraphs[0]
            p.text = val
            p.font.name = FONT_FAMILY
            if r_idx == 0:
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = AWS_ORANGE if c_idx == 0 else ACCENT_BLUE
            else:
                p.font.size = Pt(9)
                if c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = TEXT_WHITE
                elif c_idx == 3:
                    p.font.color.rgb = ACCENT_GREEN
                else:
                    p.font.color.rgb = TEXT_MUTED

    add_footer(slide5, 5, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 6: Asynchronous CDC Pipeline & Data Dependency Flow
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "צינור עיבוד אסינכרוני (CDC) ואינדוקס חיפוש",
               "הפרדה מונעת-אירועים: DynamoDB Streams ← יומן ביקורת ב-S3 ותור SQS ← אינדקס OpenSearch",
               "ארכיטקטורת אירועים ו-CQRS")

    add_card(slide6, Inches(0.8), Inches(1.7), Inches(6.0), Inches(5.1), CARD_BG, CARD_BORDER)
    left_b = slide6.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.5), Inches(4.7))
    tf = left_b.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "שלבי מחזור החיים של עיבוד האירוע"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    steps = [
        ("1. ביצוע טרנזקציה ב-DynamoDB:", "העלאה, גרסה חדשה, עריכת מטא-דאטה או מחיקה לוגית נרשמים ב-DynamoDB ופולטים אירוע ל-DynamoDB Streams."),
        ("2. מעבד האירועים (StreamProcessor Lambda):", "קולט אצוות אירועים בזמן אמת ומבצע שתי פעולות במקביל:\n • כותב רשומת ביקורת JSON בלתי-הפיכה ל-S3 Audit Bucket.\n • שולח הודעת אינדוקס לתור ה-SQS של מנוע החיפוש."),
        ("3. מנוע האינדוקס (SQS Indexer Lambda):", "מושך הודעות מ-SQS, קורא את ה-Annotation המוסמך מ-S3 ומעדכן את אינדקס החיפוש ב-OpenSearch Serverless."),
        ("4. בידוד כשלים ותור הודעות שגויות (DLQ):", "במידה ו-OpenSearch אינו זמין זמנית, מתבצעים 3 ניסיונות חוזרים וההודעה מבודדת ב-DLQ תוך הפעלת התראת CloudWatch.")
    ]
    for title, desc in steps:
        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        p.space_before = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(2)

    add_card(slide6, Inches(7.0), Inches(1.7), Inches(5.533), Inches(5.1), CARD_BG, CARD_BORDER)
    right_b = slide6.shapes.add_textbox(Inches(7.25), Inches(1.9), Inches(5.033), Inches(4.7))
    tf_r = right_b.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "מבנה הודעת האינדוקס בתור SQS (JSON)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    p_code = tf_r.add_paragraph()
    p_code.text = (
        '{\n'
        '  "event_id": "e3b0c442-98fc-1c14-9afb-f4c8996fb924",\n'
        '  "event_type": "DOCUMENT_METADATA_UPDATED",\n'
        '  "document_id": "550e8400-e29b-41d4-a716-446655440000",\n'
        '  "document_class": "loan_agreement",\n'
        '  "application_version": 2,\n'
        '  "metadata_revision": 3,\n'
        '  "status": "ACTIVE",\n'
        '  "timestamp": "2026-08-18T15:05:00.000Z",\n'
        '  "actor_id": "usr_editor_02"\n'
        '}'
    )
    p_code.font.name = "Consolas"
    p_code.font.size = Pt(9.5)
    p_code.font.color.rgb = ACCENT_BLUE
    p_code.space_before = Pt(6)

    p_res = tf_r.add_paragraph()
    p_res.text = "יתרונות חוסן והתאוששות עצמית:"
    p_res.font.name = FONT_FAMILY
    p_res.font.size = Pt(11)
    p_res.font.bold = True
    p_res.font.color.rgb = ACCENT_PURPLE
    p_res.space_before = Pt(8)

    res_items = [
        "• צריכה אידמפוטנטית: אינדוקס מוגן גרסאות מונע דריסה של אירועים שהגיעו שלא לפי הסדר.",
        "• שחזור מלא של אינדקס החיפוש: ניתן למחוק ולבנות מחדש 100% מאינדקס OpenSearch מתוך נתוני ה-S3 Annotations.",
        "• בידוד תקלות: השבתה של מנוע החיפוש אינה מפריעה להעלאות, הורדות או עריכת מטא-דאטה."
    ]
    for item in res_items:
        p = tf_r.add_paragraph()
        p.text = item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(2)

    add_footer(slide6, 6, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 7: Triple-Level Version Lineage & Field Mutability Matrix
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "מעקב גרסאות תלת-שכבתי ומטריצת שינוי שדות",
               "שלמות קריפטוגרפית, מעקב מצב מונוטוני ומשילות סכמות קפדנית ללא אובדן מידע",
               "ניהול גרסאות ומשילות מטא-דאטה")

    counters = [
        ("1. גרסת יישום לוגית", "application_version (1, 2, 3...)", "מספר שלם מונוטוני העולה אך ורק בהעלאת קובץ בינארי חדש (POST /versions). מאפס את מונה ה-revision ל-1.", ACCENT_BLUE),
        ("2. גרסה פיזית ב-S3", "s3_version_id (\"3/L4bqt9...\")", "מזהה גרסה קריפטוגרפי ובלתי-הפיך הנוצר אוטומטית ע\"י S3 עבור הקובץ הבינארי המאוחסן ב-Bucket הראשי.", ACCENT_GREEN),
        ("3. גרסת מטא-דאטה (OCC)", "metadata_revision (1, 2, 3...)", "מונה עדכוני מטא-דאטה אופטימי העולה בכל עדכון (PATCH /metadata). נאכף אטומית ב-DynamoDB.", ACCENT_AMBER)
    ]
    cw = Inches(3.75)
    cgap = Inches(0.24)
    cx_start = Inches(0.8)
    cy = Inches(1.65)
    ch = Inches(1.5)

    for i, (name, key, desc, col) in enumerate(counters):
        cx = cx_start + i * (cw + cgap)
        add_card(slide7, cx, cy, cw, ch, CARD_BG, CARD_BORDER)

        tb = slide7.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.12), cw - Inches(0.3), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col

        p_k = tf.add_paragraph()
        p_k.text = key
        p_k.font.name = "Consolas"
        p_k.font.size = Pt(9)
        p_k.font.color.rgb = TEXT_WHITE

        db = slide7.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.65), cw - Inches(0.3), Inches(0.75))
        tf_d = db.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(8.5)
        p_d.font.color.rgb = TEXT_MUTED

    mat_y = Inches(3.35)
    mat_rows = [
        ("שמות שדות במערכת", "קטגוריה", "כלל שינוי (Mutability)", "מנגנון עדכון ומשילות"),
        ("document_id, document_class, created_at, created_by", "זהות מערכתית", "בלתי-הפיך לחלוטין", "מוקצה פעם אחת ביצירת המסמך. ניסיון שינוי נדחה בשגיאת 400 Bad Request."),
        ("application_version, content_type, content_length, content_checksum", "מפרט בינארי", "בלתי-הפיך עבור הגרסה", "צמוד קריפטוגרפית לקובץ (SHA-256). מתעדכן אך ורק בקריאת POST /versions."),
        ("metadata_revision, metadata_updated_at, metadata_updated_by", "מעקב מקביליות", "מונוטוני אוטומטי", "מנוהל אוטומטית ע\"י פונקציית ה-Lambda ומנגנון ה-OCC ב-DynamoDB."),
        ("customer_id, complete_customer_id_code, account_id, business_area_code, loan_number...", "תכונות עסקיות ומשותפות", "ניתן לעריכה (PATCH)", "ניתן לעדכון ב-PATCH /metadata עם אימות סכמת JSON Schema ובדיקת מונה גרסה.")
    ]

    mat_table_shape = slide7.shapes.add_table(len(mat_rows), 4, Inches(0.8), mat_y, Inches(11.733), Inches(3.45))
    mat_table = mat_table_shape.table
    mat_table.columns[0].width = Inches(3.6)
    mat_table.columns[1].width = Inches(2.0)
    mat_table.columns[2].width = Inches(2.3)
    mat_table.columns[3].width = Inches(3.833)

    for r_idx, row in enumerate(mat_rows):
        for c_idx, val in enumerate(row):
            cell = mat_table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(22, 30, 46)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.1)
            tf.margin_top = tf.margin_bottom = Inches(0.06)

            p = tf.paragraphs[0]
            p.text = val
            p.font.name = FONT_FAMILY
            if r_idx == 0:
                p.font.size = Pt(9.5)
                p.font.bold = True
                p.font.color.rgb = AWS_ORANGE if c_idx == 0 else ACCENT_BLUE
            else:
                p.font.size = Pt(8.5)
                if c_idx == 2:
                    p.font.bold = True
                    if "לחלוטין" in val:
                        p.font.color.rgb = ACCENT_ROSE
                    elif "עבור הגרסה" in val:
                        p.font.color.rgb = ACCENT_AMBER
                    elif "מונוטוני" in val:
                        p.font.color.rgb = ACCENT_PURPLE
                    else:
                        p.font.color.rgb = ACCENT_GREEN
                elif c_idx == 0:
                    p.font.name = "Consolas"
                    p.font.color.rgb = TEXT_WHITE
                else:
                    p.font.color.rgb = TEXT_MUTED

    add_footer(slide7, 7, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 8: Metadata Lifecycle & Optimistic Concurrency Control (OCC)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "בקרת מקביליות אופטימית (OCC) והתפתחות סכמות",
               "מניעת דריסת עדכונים בעומסים גבוהים ואימות סכמות JSON Schema מהיר בזיכרון (Ajv)",
               "בקרת מקביליות ומשילות")

    add_card(slide8, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    left_b = slide8.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf = left_b.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "אלגוריתם ה-OCC (PATCH /v1/documents/{id}/metadata)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    occ_steps = [
        ("1. הלקוח מעביר גרסה צפויה:", "הקריאה כוללת expected_metadata_revision: N יחד עם השדות המבוקשים לעדכון."),
        ("2. בדיקת עקביות ב-DynamoDB:", "ה-Lambda קורא את רשומת DOC#{id}. אם המונה אינו תואם בדיוק ל-N, הקריאה נדחית מיד בשגיאת 409 Conflict."),
        ("3. שליפת ה-Annotation מ-S3 ומיזוג:", "ה-Lambda קורא את ה-JSON הקיים מ-S3, מוודא שאין שינוי בשדות אסורים, מעלה את המונה ל-N+1 ומאמת מול Ajv."),
        ("4. שמירת ה-Annotation המעודכן ב-S3:", "שמירה אטומית של ה-JSON המעודכן ישירות ב-S3 באמצעות PutObjectAnnotation."),
        ("5. עדכון מותנה אטומי ב-DynamoDB:", "ביצוע עדכון מותנה (ConditionExpression: current_metadata_revision = :expRev). במידה והתרחש מרוץ מקבילי, הפעולה נכשלת.")
    ]
    for title, desc in occ_steps:
        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        p.space_before = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(1)

    add_card(slide8, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    right_b = slide8.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf_r = right_b.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "משילות סכמות ואימות בשיטת GitOps"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    p_sch = tf_r.add_paragraph()
    p_sch.text = "מחזור חיים מבוסס GitOps (אימות בזמן הידור):"
    p_sch.font.name = FONT_FAMILY
    p_sch.font.size = Pt(10.5)
    p_sch.font.bold = True
    p_sch.font.color.rgb = ACCENT_PURPLE
    p_sch.space_before = Pt(6)

    schema_points = [
        ("• אפס השהיית Cold Start: ", "סכמות ה-JSON נארזות ישירות בקוד ה-Lambda ומהודרות לזיכרון בעלייה (Ajv Precompiled) ללא צורך בקריאות רשת."),
        ("• הגדרת מעטפת סגורה: ", "הגדרת additionalProperties: false מונעת זיהום נתונים ושינויי מבנה לא מבוקרים."),
        ("• בקרת שינויים קפדנית ב-Git: ", "כל עדכון סכמה דורש Pull Request, בדיקות linting אוטומטיות ותיעוד מלא."),
        ("• שלב 2 במפת הדרכים: ", "תמיכה עתידית ב-Schema Registry דינמי מבוסס DynamoDB להרשמת סכמות בזמן ריצה.")
    ]
    for label, text in schema_points:
        p = tf_r.add_paragraph()
        p.text = label
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_before = Pt(5)

        run = p.add_run()
        run.text = text
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    add_footer(slide8, 8, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 9: Bimodal Ingestion Architecture
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "ארכיטקטורת קליטה דו-ערוצית (Bimodal Ingestion)",
               "אופטימיזציית זמני תגובה וביטול מגבלת 10 MiB של API Gateway עבור קבצים גדולים",
               "קליטת מסמכים מתקדמת")

    add_card(slide9, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    left_b = slide9.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf = left_b.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "ערוץ 1: העלאה ישירה בבקשה אחת (≤ 4 MiB)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    p_sub = tf.add_paragraph()
    p_sub.text = "בקשת HTTP בודדת למסמכים רגילים, קבלות, טפסי הלוואה ותעודות זהות."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9.5)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(4)

    p_box = tf.add_paragraph()
    p_box.text = (
        'POST /v1/documents\n'
        'Headers:\n'
        '  • Content-Type: application/pdf\n'
        '  • X-Content-SHA256: {hex_hash}\n'
        '  • X-Document-Metadata: {base64_json_envelope}\n'
        'Body: Raw Binary Bytes'
    )
    p_box.font.name = "Consolas"
    p_box.font.size = Pt(9)
    p_box.font.color.rgb = ACCENT_GREEN
    p_box.space_before = Pt(8)

    p_adv = tf.add_paragraph()
    p_adv.text = "מאפייני ביצועים:"
    p_adv.font.name = FONT_FAMILY
    p_adv.font.size = Pt(10.5)
    p_adv.font.bold = True
    p_adv.font.color.rgb = TEXT_WHITE
    p_adv.space_before = Pt(8)

    p_adv_items = [
        "• זמן תגובה כולל: מענה סינכרוני מהיר בתוך 120ms עד 250ms בלבד.",
        "• שמירה אטומית: שמירת הקובץ ב-S3, ה-Annotation ומצביע ה-DynamoDB בטרנזקציה אחת.",
        "• חוויית לקוח פשוטה: קריאת HTTP בודדת המחזירה קוד 201 Created עם מטא-דאטה מלא."
    ]
    for item in p_adv_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(2)

    add_card(slide9, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    right_b = slide9.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf_r = right_b.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "ערוץ 2: העלאה ישירה בקישור חתום ל-S3 (> 4 MiB)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    p_sub_r = tf_r.add_paragraph()
    p_sub_r.text = "פרוטוקול דו-שלבי המיועד לתיקי משכנתאות עבי כרס, סריקות וקבצי ענק."
    p_sub_r.font.name = FONT_FAMILY
    p_sub_r.font.size = Pt(9.5)
    p_sub_r.font.color.rgb = TEXT_MUTED
    p_sub_r.space_before = Pt(4)

    p_box_r = tf_r.add_paragraph()
    p_box_r.text = (
        'שלב 1: POST /v1/documents/uploads\n'
        '  ← אימות סכמה ויצירת קישור S3 Presigned PUT ל-15 דקות.\n'
        'שלב 2: PUT {presigned_url}\n'
        '  ← הלקוח מזרים את הקובץ ישירות ל-S3 במהירות מקסימלית.\n'
        'שלב 3: POST /v1/uploads/{id}/complete\n'
        '  ← ה-Lambda מוודא תקינות, כותב Annotation ומפעיל את המסמך.'
    )
    p_box_r.font.name = "Consolas"
    p_box_r.font.size = Pt(9)
    p_box_r.font.color.rgb = ACCENT_BLUE
    p_box_r.space_before = Pt(8)

    p_adv_r = tf_r.add_paragraph()
    p_adv_r.text = "מאפייני ביצועים:"
    p_adv_r.font.name = FONT_FAMILY
    p_adv_r.font.size = Pt(10.5)
    p_adv_r.font.bold = True
    p_adv_r.font.color.rgb = TEXT_WHITE
    p_adv_r.space_before = Pt(8)

    p_adv_items_r = [
        "• עקיפת מגבלת 10 MiB ב-API Gateway: תמיכה בקבצים עד גודל של עשרות גיגה-בייט.",
        "• אפס תקורת זיכרון ב-Lambda: שכבת המחשוב אינה מתווכת בנפח הקבצים עצמם.",
        "• ניקוי אוטומטי מבוסס TTL: סשנים שלא הושלמו מנוקים אוטומטית לאחר 24 שעות ב-DynamoDB."
    ]
    for item in p_adv_items_r:
        p = tf_r.add_paragraph()
        p.text = item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(2)

    add_footer(slide9, 9, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 10: Exposed REST APIs Comprehensive Catalog
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "קטלוג ה-REST APIs המלא: 15 נקודות קצה ייעודיות",
               "חלוקה פונקציונלית לקליטה, מחזור חיים, מטא-דאטה, הורדה מאובטחת וחיפוש מתקדם",
               "ממשקי API וחוזים")

    api_catalog = [
        ("מתודה ונתיב", "שם הפעולה", "רכיבי יעד מעורבים", "תפקיד RBAC נדרש", "סוג"),
        ("GET /health", "בדיקת בריאות מערכתית", "DynamoDB, S3, OpenSearch", "פתוח / כל התפקידים", "סינכרוני"),
        ("POST /v1/documents", "העלאת מסמך ישירה (Inline)", "API GW, Lambda, S3, DynamoDB", "Document.Writer / Admin", "סינכרוני"),
        ("POST /v1/documents/uploads", "אתחול העלאה ישירה ל-S3", "Lambda, S3 Presign, DynamoDB", "Document.Writer / Admin", "סינכרוני"),
        ("POST /v1/uploads/{id}/complete", "השלמת העלאה ישירה", "Lambda, S3 Annotate, DynamoDB", "Document.Writer / Admin", "סינכרוני"),
        ("DELETE /v1/uploads/{id}", "ביטול סשן העלאה ישירה", "Lambda, S3 Abort, DynamoDB", "Document.Writer / Admin", "סינכרוני"),
        ("GET /v1/documents/{id}", "שליפת מצביע מסמך", "API GW, Lambda, DynamoDB", "Document.Reader / כולם", "סינכרוני"),
        ("GET /v1/documents/{id}/versions", "שליפת היסטוריית גרסאות", "Lambda, DynamoDB VER# Query", "Document.Reader / כולם", "סינכרוני"),
        ("POST /v1/documents/{id}/versions", "העלאת גרסה בינארית חדשה", "Lambda, S3 Binary+Annot, DynamoDB", "Document.Writer / Admin", "סינכרוני"),
        ("GET /v1/documents/{id}/versions/{v}", "שליפת גרסה היסטורית", "Lambda, DynamoDB VER# Item", "Document.Reader / כולם", "סינכרוני"),
        ("GET /v1/documents/{id}/metadata", "שליפת מטא-דאטה מוסמך", "Lambda, S3 Annotation Get", "Document.Reader / כולם", "סינכרוני"),
        ("PATCH /v1/documents/{id}/metadata", "עדכון מטא-דאטה (OCC)", "Lambda, S3 Annotate, DynamoDB", "Document.MetadataEditor / Writer", "סינכרוני"),
        ("GET /v1/documents/{id}/download", "קבלת קישור הורדה חתום", "Lambda, S3 Presign (15 min)", "Document.Reader / כולם", "סינכרוני"),
        ("POST /v1/documents/{id}/soft-delete", "מחיקה לוגית (Soft Delete)", "Lambda, DynamoDB, CDC De-index", "Document.Admin בלבד", "סינכרוני"),
        ("POST /v1/documents/{id}/restore", "שחזור מסמך שנמחק לוגית", "Lambda, DynamoDB, CDC Re-index", "Document.Admin בלבד", "סינכרוני"),
        ("POST /v1/search", "חיפוש טקסטואלי וסינון", "Lambda, OpenSearch Serverless", "Document.Reader / כולם", "סינכרוני")
    ]

    api_table_shape = slide10.shapes.add_table(len(api_catalog), 5, Inches(0.8), Inches(1.65), Inches(11.733), Inches(5.1))
    api_table = api_table_shape.table
    api_table.columns[0].width = Inches(2.6)
    api_table.columns[1].width = Inches(2.3)
    api_table.columns[2].width = Inches(3.1)
    api_table.columns[3].width = Inches(2.7)
    api_table.columns[4].width = Inches(1.033)

    for r_idx, row in enumerate(api_catalog):
        for c_idx, val in enumerate(row):
            cell = api_table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(22, 30, 46)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.04)

            p = tf.paragraphs[0]
            p.text = val
            p.font.name = FONT_FAMILY
            if r_idx == 0:
                p.font.size = Pt(9.5)
                p.font.bold = True
                p.font.color.rgb = AWS_ORANGE if c_idx == 0 else ACCENT_BLUE
            else:
                p.font.size = Pt(8.5)
                if c_idx == 0:
                    p.font.name = "Consolas"
                    if "POST" in val:
                        p.font.color.rgb = ACCENT_GREEN
                    elif "GET" in val:
                        p.font.color.rgb = ACCENT_BLUE
                    elif "PATCH" in val:
                        p.font.color.rgb = ACCENT_AMBER
                    elif "DELETE" in val:
                        p.font.color.rgb = ACCENT_ROSE
                    else:
                        p.font.color.rgb = TEXT_WHITE
                elif c_idx == 3:
                    if "Admin" in val:
                        p.font.color.rgb = ACCENT_PURPLE
                    else:
                        p.font.color.rgb = TEXT_WHITE
                else:
                    p.font.color.rgb = TEXT_MUTED

    add_footer(slide10, 10, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 11: Sequence Deep Dive - Inline Binary Upload (POST /v1/documents)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "פירוט תהליך רצף: העלאת מסמך ישירה (Inline ≤ 4 MiB)",
               "תהליך רצף 02: POST /v1/documents — שמירה אטומית בטרנזקציה אחת",
               "רצפי הרצת ממשקים (API Sequences)")

    img_02 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/02_post_documents_inline.jpg")
    if os.path.exists(img_02):
        slide11.shapes.add_picture(img_02, Inches(0.8), Inches(1.65), Inches(7.8), Inches(5.15))

    add_card(slide11, Inches(8.8), Inches(1.65), Inches(3.733), Inches(5.15), CARD_BG, CARD_BORDER)
    sb11 = slide11.shapes.add_textbox(Inches(9.0), Inches(1.85), Inches(3.333), Inches(4.75))
    tf = sb11.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "פירוט שלבי הביצוע"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    points11 = [
        ("1. בקרת כניסה ואימות:", "API Gateway מוודא את טוקן ה-JWT מול Cognito (דורש הרשאת Document.Writer או Admin)."),
        ("2. אימות סכמה וחישוב Checksum:", "ה-Lambda מחשב SHA-256 על הקובץ ומאמת את מעטפת ה-JSON מול הסכמה בזיכרון (Ajv)."),
        ("3. כתיבה ל-S3 והצמדת Annotation:", "העלאת הקובץ ל-S3 והצמדת ה-Annotation המובנה (document-metadata)."),
        ("4. שמירה אטומית ב-DynamoDB:", "ביצוע TransactWriteItems ליצירת מצביע DOC#{id}, גרסה VER#0000000001 ונעילת Idempotency."),
        ("5. פליטת אירוע CDC:", "זרמי DynamoDB מפעילים אסינכרונית את מעבד האירועים לרישום ביקורת ואינדוקס ב-OpenSearch.")
    ]
    for title, desc in points11:
        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_before = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(1)

    add_footer(slide11, 11, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 12: Sequence Deep Dive - Direct Upload Initiation & Complete
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_header(slide12, "פירוט תהליך רצף: העלאה ישירה בקישור חתום ל-S3",
               "תהליכי רצף 03 ו-04: POST /v1/documents/uploads ו-POST /v1/uploads/{id}/complete",
               "רצפי הרצת ממשקים (API Sequences)")

    img_03 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/03_post_documents_uploads_init.jpg")
    img_04 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/04_post_uploads_complete.jpg")

    if os.path.exists(img_03):
        slide12.shapes.add_picture(img_03, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_04):
        slide12.shapes.add_picture(img_04, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide12, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b12 = slide12.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf12 = b12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "עיקרי פרוטוקול ההעלאה הדו-שלבי:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    p_sub = tf12.add_paragraph()
    p_sub.text = "• שלב 1 (אתחול - 03): אימות סכמת מטא-דאטה, יצירת סשן UPLOAD#{id} ב-DynamoDB עם TTL ל-24 שעות, והנפקת קישור חתום S3 Presigned PUT ל-15 דקות.\n• העברה ישירה: הלקוח מזרים את הקובץ ישירות ל-S3 ללא מעבר ב-API Gateway או Lambda.\n• שלב 2 (השלמה - 04): ה-Lambda מאמת את הקובץ שהועלה (HeadObject), כותב את ה-Annotation ל-S3, מעדכן את הסשן ומפעיל את מצביע המסמך ב-DynamoDB."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide12, 12, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 13: Sequence Deep Dive - Upload Abort & Document Retrieval
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13)
    add_header(slide13, "פירוט תהליך רצף: ביטול העלאה ושליפת פרטי מסמך",
               "תהליכי רצף 05 ו-06: DELETE /v1/uploads/{id} ו-GET /v1/documents/{id}",
               "רצפי הרצת ממשקים (API Sequences)")

    img_05 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/05_delete_uploads_cancel.jpg")
    img_06 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/06_get_documents_document_id.jpg")

    if os.path.exists(img_05):
        slide13.shapes.add_picture(img_05, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_06):
        slide13.shapes.add_picture(img_06, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide13, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b13 = slide13.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf13 = b13.text_frame
    tf13.word_wrap = True

    p = tf13.paragraphs[0]
    p.text = "עיקרי תהליכי הביטול והשליפה המהירה:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ROSE

    p_sub = tf13.add_paragraph()
    p_sub.text = "• ביטול העלאה (05): מעדכן את סשן ההעלאה ב-DynamoDB לסטטוס ABORTED ומנקה קבצי Staging שלא הושלמו מ-S3 למניעת צבירת קבצים יתומים.\n• שליפת מצביע מסמך (06): קריאה אולטרה-מהירה O(1) בעקביות מלאה מ-DynamoDB. מחזירה מספר גרסה פעילה, מזהה גרסת S3, סטטוס ומונה עדכוני מטא-דאטה בפחות מ-10ms."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide13, 13, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 14: Sequence Deep Dive - Version Lineage & Binary Updates
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14)
    add_header(slide14, "פירוט תהליך רצף: היסטוריית גרסאות ועדכון קובץ בינארי",
               "תהליכי רצף 07 ו-08: GET /v1/documents/{id}/versions ו-POST /v1/documents/{id}/versions",
               "רצפי הרצת ממשקים (API Sequences)")

    img_07 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/07_get_documents_versions.jpg")
    img_08 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/08_post_documents_versions.jpg")

    if os.path.exists(img_07):
        slide14.shapes.add_picture(img_07, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_08):
        slide14.shapes.add_picture(img_08, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide14, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b14 = slide14.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf14 = b14.text_frame
    tf14.word_wrap = True

    p = tf14.paragraphs[0]
    p.text = "מנגנון שרשרת הגרסאות ועדכון תוכן בינארי:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE

    p_sub = tf14.add_paragraph()
    p_sub.text = "• שליפת היסטוריית גרסאות (07): שאילתת DynamoDB מהירה (pk = DOC#{id} AND begins_with(sk, 'VER#')) המחזירה את כל היסטוריית הגרסאות ללא צורך ב-ListObjects איטי ב-S3.\n• העלאת גרסה בינארית חדשה (08): העלאת קובץ חדש לאותו מפתח S3, יצירת VersionId חדש, קידום application_version (למשל 1 ← 2), איפוס מונה המטא-דאטה ל-1, ורישום רשומת VER# חדשה ב-DynamoDB."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide14, 14, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 15: Sequence Deep Dive - Historical Version Details & S3 Metadata
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15)
    add_header(slide15, "פירוט תהליך רצף: גרסאות היסטוריות ושליפת מטא-דאטה מוסמך",
               "תהליכי רצף 09 ו-10: GET /v1/documents/{id}/versions/{v} ו-GET /v1/documents/{id}/metadata",
               "רצפי הרצת ממשקים (API Sequences)")

    img_09 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/09_get_documents_version_id.jpg")
    img_10 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/10_get_documents_metadata.jpg")

    if os.path.exists(img_09):
        slide15.shapes.add_picture(img_09, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_10):
        slide15.shapes.add_picture(img_10, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide15, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b15 = slide15.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf15 = b15.text_frame
    tf15.word_wrap = True

    p = tf15.paragraphs[0]
    p.text = "שליפת מידע היסטורי ומטא-דאטה מוסמך מ-S3:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    p_sub = tf15.add_paragraph()
    p_sub.text = "• פרטי גרסה היסטורית (09): שליפת רשומת הגרסה ההיסטורית (pk = DOC#{id}, sk = VER#{padded_version}) מ-DynamoDB, הכוללת גודל קובץ, Checksum ומזהה S3 VersionId המקורי.\n• שליפת מטא-דאטה מוסמך (10): איתור מצביע הגרסה הפעילה מ-DynamoDB ושליפת ה-Annotation המלא והמוסמך ישירות מגרסת אובייקט ה-S3 באמצעות GetObjectAnnotation."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide15, 15, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 16: Sequence Deep Dive - Optimistic Metadata Patch (API 11)
    # =========================================================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16)
    add_header(slide16, "פירוט תהליך רצף: עדכון מטא-דאטה בבקרת מקביליות אופטימית",
               "תהליך רצף 11: PATCH /v1/documents/{id}/metadata — מניעת מרוצים ודריסות",
               "רצפי הרצת ממשקים (API Sequences)")

    img_11 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/11_patch_documents_metadata.jpg")
    if os.path.exists(img_11):
        slide16.shapes.add_picture(img_11, Inches(0.8), Inches(1.65), Inches(7.8), Inches(5.15))

    add_card(slide16, Inches(8.8), Inches(1.65), Inches(3.733), Inches(5.15), CARD_BG, CARD_BORDER)
    sb16 = slide16.shapes.add_textbox(Inches(9.0), Inches(1.85), Inches(3.333), Inches(4.75))
    tf = sb16.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "שלבי תהליך העדכון (OCC)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER

    points16 = [
        ("1. בדיקת גרסה ב-DynamoDB:", "קריאה עקבית של DOC#{id} והשוואת expected_metadata_revision למונה הנוכחי (דחייה ב-409 במקרה של אי-התאמה)."),
        ("2. שליפת ה-Annotation הקיים:", "קריאת מסמך ה-JSON הקיים מגרסת ה-S3 באמצעות GetObjectAnnotation."),
        ("3. אימות סכמה וערכים אסורים:", "דחיית כל ניסיון שינוי של שדות מערכתיים בלתי-הפיכים, קידום המונה ל-N+1 ואימות מלא מול Ajv."),
        ("4. כתיבת ה-Annotation ל-S3:", "שמירת ה-JSON המעודכן ישירות ב-S3 באמצעות PutObjectAnnotationCommand."),
        ("5. עדכון מותנה ב-DynamoDB:", "עדכון מונה הגרסה ב-DynamoDB עם תנאי ConditionExpression המבטיח אטומיות מוחלטת.")
    ]
    for title, desc in points16:
        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_before = Pt(5)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(1)

    add_footer(slide16, 16, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 17: Sequence Deep Dive - Download, Soft Delete & Restore
    # =========================================================================
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17)
    add_header(slide17, "פירוט תהליך רצף: הורדה ישירה, מחיקה לוגית ושחזור",
               "תהליכי רצף 12, 13 ו-14: גישה ישירה לקבצים ומשילות מחזור חיים שאינה הרסנית",
               "רצפי הרצת ממשקים (API Sequences)")

    img_12 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/12_get_documents_download.jpg")
    img_13 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/13_post_documents_soft_delete.jpg")

    if os.path.exists(img_12):
        slide17.shapes.add_picture(img_12, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_13):
        slide17.shapes.add_picture(img_13, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide17, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b17 = slide17.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf17 = b17.text_frame
    tf17.word_wrap = True

    p = tf17.paragraphs[0]
    p.text = "עיקרי תהליכי ההורדה, המחיקה הלוגית והשחזור:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    p_sub = tf17.add_paragraph()
    p_sub.text = "• הורדה מאובטחת ישירה (12): הפקת קישור זמני (15 דקות) S3 Presigned GET. הלקוח מוריד ישירות מ-S3 ללא עלויות תעבורה או עומס על ה-Lambda.\n• מחיקה לוגית מבוקרת (13): עדכון סטטוס 'SOFT_DELETED' ב-DynamoDB והסרת המסמך מאינדקס OpenSearch. קבצי ה-WORM וה-Annotations ב-S3 נשמרים ללא שינוי לעמידה ברגולציה ולצורך שחזור מהיר (14)."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide17, 17, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 18: On-Demand Format Conversion & Derivative Caching (JPEG/PNG to PDF)
    # =========================================================================
    slide18 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide18)
    add_header(slide18, "המרת פורמטים לפי דרישה ושכבת נגזרות (JPEG/PNG ל-PDF)",
               "טרנספורמציה דינמית בזיכרון, שמירת נגזרות ב-S3 ושמירה על עקרון אי-ההשתנות (WORM)",
               "שכבת נגזרות קריאה (Derived Read Projections)")

    col_w = Inches(3.644)
    col_gap = Inches(0.4)
    c1_x = Inches(0.8)
    c2_x = c1_x + col_w + col_gap
    c3_x = c2_x + col_w + col_gap
    c_y = Inches(1.65)
    c_h = Inches(3.75)

    # Card 1: In-Memory Conversion
    add_card(slide18, c1_x, c_y, col_w, c_h, CARD_BG, CARD_BORDER)
    b1 = slide18.shapes.add_textbox(c1_x + Inches(0.2), c_y + Inches(0.2), col_w - Inches(0.4), c_h - Inches(0.4))
    tf1 = b1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "⚡ המרה בזיכרון (In-Memory)"
    p1.font.name = FONT_FAMILY
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_BLUE

    points1 = [
        ("מנוע TypeScript טהור: ", "שימוש בספריית pdf-lib לעיבוד מהיר בזיכרון ללא צורך בבינאריים כבדים (ImageMagick/Ghostscript)."),
        ("מעבדי AWS Graviton ARM64: ", "ביצוע בתוך פונקציות ה-Lambda הקיימות בזמן של כ-120 מילישניות בלבד."),
        ("קריאה פשוטה בממשק: ", "הפעלה שקופה באמצעות פרמטר `?format=pdf` בממשקי שליפת מסמך, גרסה והורדה.")
    ]
    for lbl, desc in points1:
        p = tf1.add_paragraph()
        p.text = lbl
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = AWS_ORANGE
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = desc
        r.font.name = FONT_FAMILY
        r.font.size = Pt(9)
        r.font.color.rgb = TEXT_MUTED

    # Card 2: Authority & Derivative Caching
    add_card(slide18, c2_x, c_y, col_w, c_h, CARD_BG, CARD_BORDER)
    b2 = slide18.shapes.add_textbox(c2_x + Inches(0.2), c_y + Inches(0.2), col_w - Inches(0.4), c_h - Inches(0.4))
    tf2 = b2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "🔒 שמירת נגזרות ומעקב שושלת"
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN

    points2 = [
        ("נתיב S3 דטרמיניסטי: ", "שמירה תחת `derivatives/{class}/{id}/{s3_version_id}.pdf` לזיהוי מיידי בקריאות חוזרות."),
        ("שמירה על קודש הקודשים (WORM): ", "המרת פורמט לעולם אינה מייצרת גרסת מסמך חדשה או משנה את מצביעי DynamoDB."),
        ("תגיות מקור במטא-דאטה: ", "תגיות S3 User Metadata (`x-amz-meta-origin-s3-version-id`) מתעדות את גרסת המקור המדויקת.")
    ]
    for lbl, desc in points2:
        p = tf2.add_paragraph()
        p.text = lbl
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = desc
        r.font.name = FONT_FAMILY
        r.font.size = Pt(9)
        r.font.color.rgb = TEXT_MUTED

    # Card 3: FinOps & Lifecycle
    add_card(slide18, c3_x, c_y, col_w, c_h, CARD_BG, CARD_BORDER)
    b3 = slide18.shapes.add_textbox(c3_x + Inches(0.2), c_y + Inches(0.2), col_w - Inches(0.4), c_h - Inches(0.4))
    tf3 = b3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "💰 יעילות כלכלית ומחזור חיים"
    p3.font.name = FONT_FAMILY
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_PURPLE

    points3 = [
        ("מחיקה אוטומטית כעבור 14 יום: ", "כלל S3 Lifecycle Rule מוחק נגזרות זמניות אוטומטית למניעת ניפוח עלויות אחסון."),
        ("שיעור פגיעה ב-Cache של 80%+: ", "קריאות חוזרות מקבלות ישירות קישור Presigned ל-S3 ללא צורך בחישוב מחדש."),
        ("חיסכון מעבדי Graviton: ", "ארכיטקטורת ARM64 מספקת ביצועים גבוהים יותר ועלות נמוכה ב-20% לעומת x86.")
    ]
    for lbl, desc in points3:
        p = tf3.add_paragraph()
        p.text = lbl
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_PURPLE
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = desc
        r.font.name = FONT_FAMILY
        r.font.size = Pt(9)
        r.font.color.rgb = TEXT_MUTED

    # Bottom summary card
    add_card(slide18, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    bb = slide18.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tfb = bb.text_frame
    tfb.word_wrap = True
    pb = tfb.paragraphs[0]
    pb.text = "זרימת שירות מקצה לקצה:"
    pb.font.name = FONT_FAMILY
    pb.font.size = Pt(11)
    pb.font.bold = True
    pb.font.color.rgb = AWS_ORANGE

    pb_sub = tfb.add_paragraph()
    pb_sub.text = "1. לקוח מבצע `GET /documents/{id}?format=pdf` ← 2. ה-Lambda בודקת קיום ב-Cache ← 3. במקרה של החטאה (Miss): קריאת תמונת המקור, המרה בזיכרון ב-Graviton, כתיבת נגזרת ל-S3 ← 4. החזרת קישור Presigned מאובטח להורדת ה-PDF."
    pb_sub.font.name = FONT_FAMILY
    pb_sub.font.size = Pt(9.5)
    pb_sub.font.color.rgb = TEXT_WHITE
    pb_sub.space_before = Pt(2)

    add_footer(slide18, 18, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 19: Sequence Deep Dive - OpenSearch Search & Deep Health Probe
    # =========================================================================
    slide19 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide19)
    add_header(slide19, "פירוט תהליך רצף: חיפוש טקסטואלי ובדיקת בריאות עמוקה",
               "תהליכי רצף 15 ו-01: POST /v1/search ו-GET /health",
               "רצפי הרצת ממשקים (API Sequences)")

    img_15 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/15_post_search.jpg")
    img_01 = os.path.join(script_dir, "diagrams/api_sequence_diagrams/01_get_health.jpg")

    if os.path.exists(img_15):
        slide19.shapes.add_picture(img_15, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_01):
        slide19.shapes.add_picture(img_01, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide19, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b19 = slide19.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf19 = b19.text_frame
    tf19.word_wrap = True

    p = tf19.paragraphs[0]
    p.text = "עיקרי מנוע החיפוש ובדיקות הבריאות:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE

    p_sub = tf19.add_paragraph()
    p_sub.text = "• חיפוש OpenSearch Serverless (15): שאילתות טקסט מלא, סינונים לפי שדות מטא-דאטה, תאריכים וסטטוס מסמך בזמני תגובה של מילישניות בודדות.\n• בדיקת בריאות מעמיקה (01): בודקת סינכרונית את זמינות DynamoDB, S3 ו-OpenSearch ומחזירה תמונת מצב תפעולית מלאה עבור מערכות ניטור ו-Load Balancers."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide19, 19, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 20: UI & Backend Integration Architecture
    # =========================================================================
    slide20 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide20)
    add_header(slide20, "פורטל אינטרנט Serverless וארכיטקטורת אינטגרציית לקוחות",
               "CloudFront + S3 SPA, אימות מאובטח ב-Cognito JWT וחישוב SHA-256 בצד הלקוח",
               "ארכיטקטורת ממשק משתמש ואינטגרציה")

    ui_img = os.path.join(script_dir, "diagrams/ui_backend_integration_architecture.png")
    if os.path.exists(ui_img):
        slide20.shapes.add_picture(ui_img, Inches(0.8), Inches(1.65), Inches(8.0), Inches(5.15))

    add_card(slide20, Inches(9.0), Inches(1.65), Inches(3.533), Inches(5.15), CARD_BG, CARD_BORDER)
    sb20 = slide20.shapes.add_textbox(Inches(9.2), Inches(1.85), Inches(3.133), Inches(4.75))
    tf = sb20.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "עיקרי האינטגרציה"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    ui_points = [
        ("• סימולטור תפקידים בקליק: ", "יצירת טוקנים מיידית לבדיקת הרשאות Admin, Writer, Reader וקצין ציות."),
        ("• חתימת בקשות אוטומטית: ", "החתמת כל קריאות ה-HTTP בטוקני Cognito OAuth 2.0 / OIDC JWT."),
        ("• ממשק העלאה דו-ערוצי: ", "תמיכה שקופה בהעלאה ישירה בקישור S3 או העלאה מהירה של קבצים קטנים."),
        ("• אספקת PDF לפי דרישה: ", "כפתור הורדה דינמי להמרת תמונות JPEG/PNG ל-PDF בזמן אמת."),
        ("• חקר נתונים בזמן אמת: ", "ממשק חיפוש מהיר עם סינונים דינמיים מול OpenSearch Serverless.")
    ]
    for label, desc in ui_points:
        p = tf.add_paragraph()
        p.text = label
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_before = Pt(4)

        run = p.add_run()
        run.text = desc
        run.font.name = FONT_FAMILY
        run.font.size = Pt(8.5)
        run.font.color.rgb = TEXT_MUTED

    add_footer(slide20, 20, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 21: Cost Model & Financial Sizing / RBAC Matrix
    # =========================================================================
    slide21 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide21)
    add_header(slide21, "מטריצת הרשאות RBAC, הצפנת נתונים ובקרת WORM",
               "אכיפת הרשאות קשיחה ב-Cognito, מפתחות הצפנה ייעודיים ב-KMS ומשילות רגולטורית",
               "אבטחת מידע ומשילות הרשאות")

    left_b = slide21.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.6), Inches(4.8))
    tf = left_b.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "אכיפת מטריצת הרשאות (Cognito RBAC Matrix)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    rbac_rows = [
        ("פעולה במערכת", "Reader", "Writer", "Editor", "Admin"),
        ("POST /v1/documents", "❌", "✅", "❌", "✅"),
        ("GET /v1/documents/{id}", "✅", "✅", "✅", "✅"),
        ("GET /v1/download", "✅", "✅", "✅", "✅"),
        ("POST /versions", "❌", "✅", "❌", "✅"),
        ("PATCH /metadata", "❌", "✅", "✅", "✅"),
        ("POST /soft-delete", "❌", "❌", "❌", "✅"),
        ("POST /restore", "❌", "❌", "❌", "✅"),
        ("POST /search", "✅", "✅", "✅", "✅")
    ]

    rbac_table_shape = slide21.shapes.add_table(len(rbac_rows), 5, Inches(1.0), Inches(2.25), Inches(5.6), Inches(4.3))
    rbac_table = rbac_table_shape.table
    rbac_table.columns[0].width = Inches(2.4)
    rbac_table.columns[1].width = Inches(0.8)
    rbac_table.columns[2].width = Inches(0.8)
    rbac_table.columns[3].width = Inches(0.8)
    rbac_table.columns[4].width = Inches(0.8)

    for r_idx, row in enumerate(rbac_rows):
        for c_idx, val in enumerate(row):
            cell = rbac_table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(22, 30, 46)

            tf_c = cell.text_frame
            tf_c.word_wrap = True
            tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = Inches(0.02)
            p = tf_c.paragraphs[0]
            p.text = val
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            p.font.name = FONT_FAMILY
            if r_idx == 0:
                p.font.size = Pt(8.5)
                p.font.bold = True
                p.font.color.rgb = AWS_ORANGE if c_idx == 0 else ACCENT_BLUE
            else:
                p.font.size = Pt(8)
                if c_idx == 0:
                    p.font.name = "Consolas"
                    p.font.color.rgb = TEXT_WHITE
                else:
                    p.font.color.rgb = ACCENT_GREEN if "✅" in val else ACCENT_ROSE

    add_card(slide21, Inches(7.0), Inches(1.7), Inches(5.533), Inches(5.1), CARD_BG, CARD_BORDER)
    right_b = slide21.shapes.add_textbox(Inches(7.25), Inches(1.9), Inches(5.033), Inches(4.7))
    tf_r = right_b.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "הצפנה ארגונית ובקרת WORM"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    sec_items = [
        ("מפתח הצפנה ייעודי ב-AWS KMS (CMK):", "מפתח מרכזי (alias/doc-platform-mvp) עם רוטציה שנתית אוטומטית להצפנת דליי S3, טבלאות DynamoDB, תורי SQS ולוגים."),
        ("מדיניות אי-מחיקה מוחלטת (WORM):", "חסימה מפורשת s3:DeleteObjectVersion DENY בכלל מדיניות ה-IAM. אין אפשרות לדרוס או למחוק קבצים בינאריים מקוריים."),
        ("אפס חשיפה ציבורית ב-S3:", "חסימת גישה ציבורית מלאה (S3 Block Public Access) ואכיפת TLS 1.2+ בהצפנה במעבר (In-Transit)."),
        ("תפקידי IAM מבודדים (Micro-Perimeter):", "לכל פונקציית Lambda תפקיד IAM ייעודי המגביל את פעולותיה לרכיבים הנדרשים בלבד.")
    ]
    for title, desc in sec_items:
        p = tf_r.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_before = Pt(6)

        p2 = tf_r.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(1)

    add_footer(slide21, 21, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 22: Infrastructure as Code (AWS CDK) & Stack Architecture
    # =========================================================================
    slide22 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide22)
    add_header(slide22, "תשתית כקוד (IaC): 9 Stacks מודולריים ב-AWS CDK v2",
               "טופולוגיית TypeScript CloudFormation דטרמיניסטית, תיוג גלובלי ואפס תלויות מעגליות",
               "ארכיטקטורת פריסה ותשתית")

    stacks = [
        ("1. SecurityStack", "Cognito User Pool, App Client & KMS CMK Key", "מנהל את ספריית המשתמשים, קבוצות ה-RBAC ומפתח ההצפנה המרכזי.", ACCENT_BLUE),
        ("2. StorageStack", "S3 Primary Document Bucket & S3 Audit Bucket", "מגדיר גרסאות, מדיניות הצפנה, CORS ומדיניות נעילת WORM.", ACCENT_GREEN),
        ("3. SearchStack", "OpenSearch Serverless Collection (documents-v1)", "מקים את אוסף החיפוש הוקטורי, מדיניות הצפנה ואבטחת רשת.", ACCENT_PURPLE),
        ("4. MessagingStack", "Amazon SQS Indexing Queue & Dead-Letter Queue", "חוצץ אירועי CDC עם 3 ניסיונות חוזרים והתראות CloudWatch.", AWS_ORANGE),
        ("5. ControlPlaneStack", "DynamoDB Table (doc-platform-mvp-control)", "טבלה מרכזית בטכנולוגיית Single-Table עם PITR ו-Streams פעילים.", ACCENT_AMBER),
        ("6. ComputeStack", "21 פונקציות Graviton ARM64 AWS Lambda", "פורס שירותי Command, Query, Search, Stream ו-Indexer מבודדים.", ACCENT_BLUE),
        ("7. ApiStack", "Amazon API Gateway REST API & Cognito Authorizer", "מחבר נתיבים, מנגנון הרשאות, הגבלות קצב (Throttling) ו-CORS.", ACCENT_GREEN),
        ("8. ObservabilityStack", "התראות CloudWatch, לוחות בקרה ותיוג גלובלי", "פורס התראות DLQ/5xx, לוחות בקרה תפעוליים ותגיות Project/Environment.", ACCENT_ROSE)
    ]

    st_w = Inches(5.7)
    st_h = Inches(1.15)
    st_gap_x = Inches(0.333)
    st_gap_y = Inches(0.12)
    st_x1 = Inches(0.8)
    st_x2 = Inches(6.833)
    st_y_start = Inches(1.7)

    for idx, (name, res, desc, col) in enumerate(stacks):
        col_idx = idx % 2
        row_idx = idx // 2
        sx = st_x1 if col_idx == 0 else st_x2
        sy = st_y_start + row_idx * (st_h + st_gap_y)

        add_card(slide22, sx, sy, st_w, st_h, CARD_BG, CARD_BORDER)

        tb = slide22.shapes.add_textbox(sx + Inches(0.15), sy + Inches(0.1), st_w - Inches(0.3), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col

        p_r = tf.add_paragraph()
        p_r.text = f"• משאבים: {res}"
        p_r.font.name = FONT_FAMILY
        p_r.font.size = Pt(8.5)
        p_r.font.color.rgb = TEXT_WHITE

        p_d = tf.add_paragraph()
        p_d.text = f"• תפקיד: {desc}"
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(8)
        p_d.font.color.rgb = TEXT_MUTED

    add_footer(slide22, 22, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 23: Strategic Architecture Summary & Architectural Decision Records
    # =========================================================================
    slide23 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide23)
    add_header(slide23, "סיכום אסטרטגי ותיעוד החלטות ארכיטקטורה (ADRs)",
               "מדוע פתרון Serverless זה מספק שרידות מקסימלית, יעילות כלכלית ועמידה מלאה ברגולציה",
               "תמצית מנהלים וערך אסטרטגי")

    pillars = [
        ("1. אפס עלויות סרק (Zero Idle TCO)", "עלויות המערכת נגזרות ישירות מנפח השימוש בפועל. סביבות פיתוח ובדיקה עולות $0 בעת חוסר פעילות, וחוסכות מיליוני דולרים ברישוי.", ACCENT_BLUE),
        ("2. שרידות WORM בלתי-מוגבלת", "שרידות של 11 תשיעיות ב-S3 יחד עם ניהול גרסאות קשיח ומדיניות IAM Deny מבטיחים הגנה מלאה מפני תוכנות כופר ומחיקות שגויות.", ACCENT_GREEN),
        ("3. בידוד מוחלט של מרחב הכשל", "ארכיטקטורת אחסון CQRS מבוזרת מבטיחה שבנייה מחדש של אינדקס החיפוש לעולם לא תשבית העלאה, עריכה או הורדה של מסמכים.", ACCENT_PURPLE)
    ]
    pw = Inches(3.75)
    pgap = Inches(0.24)
    px_start = Inches(0.8)
    py = Inches(1.65)
    ph = Inches(1.7)

    for i, (title, desc, col) in enumerate(pillars):
        px = px_start + i * (pw + pgap)
        add_card(slide23, px, py, pw, ph, CARD_BG, CARD_BORDER)

        tb = slide23.shapes.add_textbox(px + Inches(0.15), py + Inches(0.15), pw - Inches(0.3), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col

        db = slide23.shapes.add_textbox(px + Inches(0.15), py + Inches(0.55), pw - Inches(0.3), Inches(1.05))
        tf_d = db.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_MUTED

    adr_y = Inches(3.55)
    adr_rows = [
        ("החלטה ארכיטקטונית", "הגישה שנבחרה", "חלופות שנדחו", "הנמקה אסטרטגית מרכזית"),
        ("מודל אחסון מטא-דאטה", "S3 Native Annotations", "כותרות S3 (2KB) / קבצי Sidecar", "עוקף מגבלת 2KB, מונע שכפול קבצים ומוכן לשאילתות SQL ישירות."),
        ("טבלת בקרת טרנזקציות", "DynamoDB Single-Table", "Aurora PostgreSQL / DocumentDB", "מהירות תגובה של מילישניות בודדות, בקרת מקביליות אופטימית ואפס תחזוקת שרתים."),
        ("טופולוגיית מנוע החיפוש", "OpenSearch Serverless", "Elasticsearch על EC2 / קלאסטר מנוהל", "ללא צורך בניהול גודל שרתים, התרחבות OCU אוטומטית ואימות IAM SigV4."),
        ("טרנספורמציית פורמטים", "המרת PDF לפי דרישה (14d S3 Cache)", "המרה מראש של כל הקבצים / אחסון קבוע", "ביצוע בזיכרון על Graviton ARM64, מחיקה אוטומטית כעבור 14 יום, שושלת WORM נשמרת."),
        ("אימות סכמות מטא-דאטה", "GitOps מהודר מראש (Ajv בזיכרון)", "שליפת סכמות מ-DB בעליית Lambda", "אפס תקורה בעליית פונקציה (Cold Start) ובקרת קוד קפדנית ב-Pull Requests.")
    ]

    adr_table_shape = slide23.shapes.add_table(len(adr_rows), 4, Inches(0.8), adr_y, Inches(11.733), Inches(3.25))
    adr_table = adr_table_shape.table
    adr_table.columns[0].width = Inches(2.5)
    adr_table.columns[1].width = Inches(2.8)
    adr_table.columns[2].width = Inches(2.8)
    adr_table.columns[3].width = Inches(3.633)

    for r_idx, row in enumerate(adr_rows):
        for c_idx, val in enumerate(row):
            cell = adr_table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(22, 30, 46)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.04)

            p = tf.paragraphs[0]
            p.text = val
            p.font.name = FONT_FAMILY
            if r_idx == 0:
                p.font.size = Pt(9.5)
                p.font.bold = True
                p.font.color.rgb = AWS_ORANGE if c_idx == 0 else ACCENT_BLUE
            else:
                p.font.size = Pt(8)
                if c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = TEXT_WHITE
                elif c_idx == 1:
                    p.font.color.rgb = ACCENT_GREEN
                elif c_idx == 2:
                    p.font.color.rgb = ACCENT_ROSE
                else:
                    p.font.color.rgb = TEXT_MUTED

    add_footer(slide23, 23, TOTAL_SLIDES)

    # -------------------------------------------------------------------------
    # Save Presentation
    # -------------------------------------------------------------------------
    prs.save(output_path)
    print(f"[SUCCESS] Hebrew Presentation generated successfully: {output_path}")

if __name__ == "__main__":
    build_hebrew_presentation()
