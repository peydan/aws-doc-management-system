#!/usr/bin/env python3
"""
Generate a comprehensive, Enterprise Solution PowerPoint presentation (.pptx)
for the AWS Document Management Platform.
Covers:
- End-to-End Architecture & Topology
- 5-Tier Authority & Consistency Model
- Data Dependencies: S3 Annotations vs DynamoDB vs OpenSearch Serverless
- Asynchronous CDC Streaming & S3 Audit Logging
- Triple-Level Version Lineage & Field Mutability Matrix
- Optimistic Concurrency Control (OCC) & Schema Evolution
- Bimodal Ingestion (Inline vs Direct S3 Presigned)
- All 15 Exposed APIs & Sequence Flows with High-Resolution Diagram Embeds
- Serverless Web Portal & Client Integration (CloudFront + S3 SPA)
- Security, RBAC & Cryptographic Governance (KMS CMK, WORM)
- AWS CDK v2 Infrastructure as Code (Modular Stacks)
- Resiliency, Failure Modes, DLQ & Reconciliation
- Strategic Architecture Conclusion & Production Acceptance
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -----------------------------------------------------------------------------
# Color Palette & Typography
# -----------------------------------------------------------------------------
BG_COLOR = RGBColor(11, 17, 32)       # #0B1120 Deep Midnight Navy
CARD_BG = RGBColor(30, 41, 59)        # #1E293B Slate Dark Card
CARD_BORDER = RGBColor(51, 65, 85)    # #334155 Slate Border
CARD_BG_LIGHT = RGBColor(24, 32, 47)  # #18202F

TEXT_WHITE = RGBColor(248, 250, 252)  # #F8FAFC
TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8
TEXT_DIM = RGBColor(100, 116, 139)    # #64748B

AWS_ORANGE = RGBColor(255, 153, 0)    # #FF9900 AWS Orange
ACCENT_BLUE = RGBColor(56, 189, 248)  # #38BDF8 Sky Blue
ACCENT_INDIGO = RGBColor(129, 140, 248) # #818CF8 Indigo
ACCENT_GREEN = RGBColor(52, 211, 153) # #34D399 Emerald Green
ACCENT_PURPLE = RGBColor(192, 132, 252) # #C084FC Purple
ACCENT_ROSE = RGBColor(244, 63, 94)   # #F43F5E Rose Red
ACCENT_AMBER = RGBColor(251, 191, 36) # #FBBF24 Amber

FONT_FAMILY = "Segoe UI"
FONT_FAMILY_TITLE = "Segoe UI"

def set_slide_background(slide, color=BG_COLOR):
    """Fills the slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, subtitle_text, category_tag="AWS Serverless Architecture"):
    """Adds a standardized modern header to a slide."""
    # Category Tag / Badge
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.3))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = 0
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = category_tag.upper()
    p_tag.font.name = FONT_FAMILY
    p_tag.font.size = Pt(9.5)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(11.7), Inches(0.55))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_FAMILY_TITLE
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

    # Subtitle
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.22), Inches(11.7), Inches(0.35))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.name = FONT_FAMILY
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = TEXT_MUTED

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    """Creates a rounded-corner or solid rectangular card container."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape

def add_footer(slide, current_slide, total_slides):
    """Adds a sleek footer line and pagination."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CARD_BORDER
    line.line.fill.background()

    box_l = slide.shapes.add_textbox(Inches(0.8), Inches(7.12), Inches(8.5), Inches(0.25))
    tf_l = box_l.text_frame
    tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
    p_l = tf_l.paragraphs[0]
    p_l.text = "AWS Document Management Platform • Solution Architecture & Sequence Specification"
    p_l.font.name = FONT_FAMILY
    p_l.font.size = Pt(8.5)
    p_l.font.color.rgb = TEXT_DIM

    box_r = slide.shapes.add_textbox(Inches(10.0), Inches(7.12), Inches(2.533), Inches(0.25))
    tf_r = box_r.text_frame
    tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
    p_r = tf_r.paragraphs[0]
    p_r.text = f"{current_slide} / {total_slides}"
    p_r.alignment = PP_ALIGN.RIGHT
    p_r.font.name = FONT_FAMILY
    p_r.font.size = Pt(8.5)
    p_r.font.bold = True
    p_r.font.color.rgb = ACCENT_BLUE


def build_presentation(output_path="AWS_Document_Management_Platform_Architecture.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    TOTAL_SLIDES = 22

    # =========================================================================
    # SLIDE 1: Title Slide (Executive Dark)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, BG_COLOR)

    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(1.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AWS_ORANGE
    bar.line.fill.background()

    tag_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(8.0), Inches(0.4))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AWS ENTERPRISE SOLUTION ARCHITECTURE"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(11.5), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Cloud-Native Enterprise Document Management Platform"
    p.font.name = FONT_FAMILY_TITLE
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    sub_box = slide1.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.0))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Decoupled Serverless Architecture for Immutable Binary Versioning, Native S3 Metadata Annotations, Optimistic Concurrency & CQRS Search Projections"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(15)
    p.font.color.rgb = ACCENT_BLUE

    highlights = [
        ("WORM Immutability", "S3 versioning with strict IAM Deny policies for regulatory compliance", ACCENT_BLUE),
        ("S3 Annotations", "Overcomes 2 KB header limits; zero sidecar sprawl with native JSON binding", ACCENT_GREEN),
        ("DynamoDB OCC", "Sub-10ms transactional control plane with atomic optimistic concurrency", AWS_ORANGE),
        ("CQRS Search", "OpenSearch Serverless derived read model with full blast radius isolation", ACCENT_PURPLE)
    ]
    card_w = Inches(2.78)
    card_gap = Inches(0.2)
    start_x = Inches(0.8)
    card_y = Inches(4.8)
    card_h = Inches(1.7)

    for i, (title, desc, color) in enumerate(highlights):
        cx = start_x + i * (card_w + card_gap)
        add_card(slide1, cx, card_y, card_w, card_h, CARD_BG, CARD_BORDER)

        t_box = slide1.shapes.add_textbox(cx + Inches(0.2), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.35))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color

        d_box = slide1.shapes.add_textbox(cx + Inches(0.2), card_y + Inches(0.6), card_w - Inches(0.4), Inches(0.95))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_MUTED

    add_footer(slide1, 1, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 2: Executive Summary & Enterprise Modernization Context
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Executive Summary & Enterprise Architecture Blueprint",
               "Modernizing Enterprise Content Management with Decoupled Cloud-Native Serverless Services",
               "Strategic Value & Context")

    add_card(slide2, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    left_box = slide2.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf = left_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Enterprise Architecture Core Principles"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ROSE

    items_left = [
        ("Consumption-Based Elasticity:", "100% serverless compute and on-demand database scaling ensure zero idle cost and automatic capacity expansion."),
        ("Decoupled Authority Model:", "Separating binary storage, metadata indexing, and full-text search prevents distributed locking and cascading outages."),
        ("Native S3 Object Annotations:", "Eliminates 2 KB header limits without sidecar files, enabling direct schema-validated metadata mutations."),
        ("Enterprise WORM & Auditability:", "Enforces immutable content versioning, strict optimistic concurrency control, and continuous CDC stream auditing.")
    ]
    for title, desc in items_left:
        p = tf.add_paragraph()
        p.text = f"• {title} "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(8)

    add_card(slide2, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    right_box = slide2.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf = right_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Architectural Capabilities & Business Outcomes"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    items_right = [
        ("Zero-Idle TCO:", "True serverless compute (Lambda) and storage (S3/DynamoDB/OpenSearch Serverless). Zero compute cost during idle periods."),
        ("11 9s Durability & Compliance:", "Amazon S3 raw binaries with WORM immutability; explicit s3:DeleteObjectVersion DENY policies prevent malicious deletion."),
        ("Sub-Second Latency SLAs:", "Inline uploads in ~150ms; direct presigned downloads in <50ms; sub-10ms DynamoDB transactional control plane."),
        ("CQRS Blast Radius Isolation:", "OpenSearch search index failures or re-indexing operations never impact core document downloads, uploads, or metadata edits.")
    ]
    for title, desc in items_right:
        p = tf.add_paragraph()
        p.text = f"✓ {title} "
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(8)

    add_footer(slide2, 2, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 3: End-to-End Solution Architecture Diagram
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "End-to-End Solution Architecture Blueprint",
               "Comprehensive AWS Serverless Topology, Ingress Control, Microservice Compute & Storage Layers",
               "System Topology")

    arch_img_path = "diagrams/aws_document_management_architecture.png"
    if os.path.exists(arch_img_path):
        slide3.shapes.add_picture(arch_img_path, Inches(0.8), Inches(1.65), Inches(8.2), Inches(5.15))

    add_card(slide3, Inches(9.2), Inches(1.65), Inches(3.333), Inches(5.15), CARD_BG, CARD_BORDER)
    side_box = slide3.shapes.add_textbox(Inches(9.4), Inches(1.85), Inches(2.933), Inches(4.75))
    tf = side_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Architecture Highlights"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    p_items = [
        ("API Gateway Ingress:", "Handles REST routing, Cognito JWT token validation, and binary payload passthrough (*/*).", ACCENT_BLUE),
        ("Lambda Micro-Perimeters:", "Stateless Command, Query, and Search handlers operating under strict least-privilege IAM roles.", ACCENT_GREEN),
        ("Native S3 Annotations:", "Authoritative JSON metadata bound directly to S3 object versions (document-metadata).", ACCENT_AMBER),
        ("DynamoDB Control Table:", "Single-table design tracking active pointers, OCC revisions, and upload sessions.", ACCENT_PURPLE),
        ("Async CDC Pipeline:", "DynamoDB Streams feed S3 Audit Bucket & SQS Indexing Queue to hydrate OpenSearch Serverless.", ACCENT_BLUE)
    ]
    for title, desc, color in p_items:
        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_before = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(9)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(1)

    add_footer(slide3, 3, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 4: The 5-Tier Authority & Consistency Model
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "The 5-Tier Authority & Consistency Model",
               "Eliminating Distributed 2-Phase Transactions by Assigning Strict Authority Boundaries",
               "Core Data Governance")

    tiers = [
        ("1. Content Authority", "Amazon S3 (Primary Bucket)", "Immutable raw binary storage (PDF, TIFF, Office). Versioned objects with strict s3:DeleteObjectVersion DENY policies for absolute WORM compliance. Raw bytes are cryptographically locked via SHA-256.", ACCENT_BLUE),
        ("2. Metadata Authority", "Amazon S3 Annotations", "Full structured JSON metadata stored as authoritative named annotations (document-metadata) bound directly to S3 object versions. Conforms strictly to schema bank.document-metadata/1.", ACCENT_GREEN),
        ("3. Control Plane Authority", "Amazon DynamoDB", "Single-table design tracking active document pointers (DOC#{id}), version lineages (VER#{pad}), optimistic concurrency locks, and upload sessions at sub-10ms latency.", AWS_ORANGE),
        ("4. Search Read Projection", "OpenSearch Serverless", "Derived CQRS read model (documents-v1) hydrated asynchronously via DynamoDB Streams and SQS. Document ingestion and downloads remain 100% operational if search is offline.", ACCENT_PURPLE),
        ("5. Non-Repudiation Audit", "Amazon S3 Audit Bucket", "Immutable CDC mutation stream logs automatically written to a KMS-encrypted, date-partitioned audit bucket (audit/{date}/{doc_id}_{event_id}.json) for forensics and compliance.", ACCENT_ROSE)
    ]

    tier_h = Inches(0.95)
    tier_gap = Inches(0.12)
    start_y = Inches(1.65)

    for i, (title, service, desc, color) in enumerate(tiers):
        ty = start_y + i * (tier_h + tier_gap)
        add_card(slide4, Inches(0.8), ty, Inches(11.733), tier_h, CARD_BG, CARD_BORDER)

        tag_b = slide4.shapes.add_textbox(Inches(1.05), ty + Inches(0.15), Inches(3.2), Inches(0.65))
        tf = tag_b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(12.5)
        p.font.bold = True
        p.font.color.rgb = color

        p_srv = tf.add_paragraph()
        p_srv.text = service
        p_srv.font.name = FONT_FAMILY
        p_srv.font.size = Pt(10)
        p_srv.font.color.rgb = TEXT_MUTED

        desc_b = slide4.shapes.add_textbox(Inches(4.4), ty + Inches(0.15), Inches(7.9), Inches(0.65))
        tf_d = desc_b.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_WHITE

    add_footer(slide4, 4, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 5: Interplay: S3 Annotations vs. DynamoDB Control Plane
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Storage Interplay: S3 Annotations vs. DynamoDB Control Table",
               "Why Both Are Essential: Resolving Object Storage Limitations with Transactional Database Speed",
               "Data Architecture Deep Dive")

    rows = [
        ("Dimension / Capability", "Amazon S3 & S3 Annotations", "Amazon DynamoDB Control Table", "Why DynamoDB Is Essential"),
        ("Authoritative Metadata", "Primary Authority (Holds full JSON schema payload)", "Stores lightweight pointer, revision count & ETag", "S3 Annotations bind rich data to object versions; enables Iceberg SQL queries."),
        ("Optimistic Concurrency (OCC)", "No conditional expressions on arbitrary JSON attributes", "Primary: Atomic ConditionExpression on revision", "DynamoDB atomically verifies current_metadata_revision = :exp to prevent lost updates."),
        ("Active Pointer Resolution", "Requires listing all object versions in bucket", "Primary: O(1) sub-millisecond key-value lookup", "DynamoDB maps DOC#{id} to latest active s3_version_id without slow S3 API listings."),
        ("Version History Lineage", "Non-sequential, opaque S3 version string IDs", "Primary: Sort-Key indexed query begins_with(VER#)", "Enables instant retrieval of ordered history (v1, v2, v3) in a single fast read query."),
        ("Multi-Item ACID Transactions", "No multi-object transactional commit API", "Primary: TransactWriteItems binds pointer & locks", "Atomically updates pointer, records version item, and commits idempotency lock in 1 call."),
        ("Change Data Capture (CDC)", "S3 Event Notifications lack before/after diffs", "Primary: Ordered DynamoDB Streams (Old/New)", "Emits rich event streams to asynchronously drive the S3 audit logger and search index.")
    ]

    table_shape = slide5.shapes.add_table(len(rows), 4, Inches(0.8), Inches(1.65), Inches(11.733), Inches(5.1))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.0)
    table.columns[2].width = Inches(3.2)
    table.columns[3].width = Inches(3.333)

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(22, 30, 46)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.12)
            tf.margin_top = tf.margin_bottom = Inches(0.08)

            p = tf.paragraphs[0]
            p.text = val
            p.font.name = FONT_FAMILY
            if r_idx == 0:
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = AWS_ORANGE if c_idx == 0 else ACCENT_BLUE
            else:
                p.font.size = Pt(9)
                if c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = TEXT_WHITE
                elif c_idx == 3:
                    p.font.color.rgb = ACCENT_GREEN
                else:
                    p.font.color.rgb = TEXT_MUTED

    add_footer(slide5, 5, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 6: Asynchronous CDC Pipeline & Data Dependency Flow
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Asynchronous CDC & Search Hydration Pipeline",
               "Event-Driven Downstream Decoupling: DynamoDB Streams → S3 Audit Bucket & SQS → OpenSearch",
               "Event Architecture & CQRS")

    add_card(slide6, Inches(0.8), Inches(1.7), Inches(6.0), Inches(5.1), CARD_BG, CARD_BORDER)
    left_b = slide6.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.5), Inches(4.7))
    tf = left_b.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Pipeline Execution Lifecycle"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    steps = [
        ("1. Transactional Commit in DynamoDB:", "Upload, versioning, metadata edit, or soft-delete commits to DynamoDB and emits a CDC record to DynamoDB Streams."),
        ("2. Stream Processor Lambda:", "Consumes batches of stream records in real time and triggers two parallel operations:\n • Writes an immutable JSON audit log to the S3 Audit Bucket.\n • Dispatches an indexing message to the Amazon SQS Indexing Queue."),
        ("3. SQS Indexer Consumer Lambda:", "Pulls batches from SQS, reads the authoritative S3 annotation, and updates the OpenSearch Serverless collection (documents-v1)."),
        ("4. Poison Message Quarantine & DLQ:", "If OpenSearch is temporarily unavailable, SQS retries 3 times before isolating failed records in doc-platform-mvp-index-dlq and triggering CloudWatch alarms.")
    ]
    for title, desc in steps:
        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        p.space_before = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(2)

    add_card(slide6, Inches(7.0), Inches(1.7), Inches(5.533), Inches(5.1), CARD_BG, CARD_BORDER)
    right_b = slide6.shapes.add_textbox(Inches(7.25), Inches(1.9), Inches(5.033), Inches(4.7))
    tf_r = right_b.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "SQS Index Message Contract (JSON)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    p_code = tf_r.add_paragraph()
    p_code.text = (
        '{\n'
        '  "event_id": "e3b0c442-98fc-1c14-9afb-f4c8996fb924",\n'
        '  "event_type": "DOCUMENT_METADATA_UPDATED",\n'
        '  "document_id": "550e8400-e29b-41d4-a716-446655440000",\n'
        '  "document_class": "loan_agreement",\n'
        '  "application_version": 2,\n'
        '  "metadata_revision": 3,\n'
        '  "status": "ACTIVE",\n'
        '  "timestamp": "2026-08-18T15:05:00.000Z",\n'
        '  "actor_id": "usr_editor_02"\n'
        '}'
    )
    p_code.font.name = "Consolas"
    p_code.font.size = Pt(9.5)
    p_code.font.color.rgb = ACCENT_BLUE
    p_code.space_before = Pt(6)

    p_res = tf_r.add_paragraph()
    p_res.text = "Resilience & Self-Healing Properties:"
    p_res.font.name = FONT_FAMILY
    p_res.font.size = Pt(11)
    p_res.font.bold = True
    p_res.font.color.rgb = ACCENT_PURPLE
    p_res.space_before = Pt(8)

    res_items = [
        "• Idempotent Consumer: Version-guarded upserts prevent out-of-order race conditions.",
        "• Index Rehydration: Search index can be dropped and 100% rebuilt from S3 annotations.",
        "• Blast Radius Isolation: Search outages never block transactional uploads or downloads."
    ]
    for item in res_items:
        p = tf_r.add_paragraph()
        p.text = item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(2)

    add_footer(slide6, 6, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 7: Triple-Level Version Lineage & Field Mutability Matrix
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Triple-Level Version Lineage & Immutability Matrix",
               "Cryptographic Integrity, Monotonic State Tracking, and Strict Schema Governance",
               "Versioning & Metadata Governance")

    counters = [
        ("1. Application Version", "application_version (1, 2, 3...)", "Monotonic integer. Increments strictly when new binary content is uploaded via POST /versions. Resets metadata_revision to 1.", ACCENT_BLUE),
        ("2. Physical S3 VersionId", "s3_version_id (\"3/L4bqt9...\")", "Immutable S3-generated cryptographic object version identifier representing raw binary bytes stored in the primary S3 bucket.", ACCENT_GREEN),
        ("3. Metadata Revision", "metadata_revision (1, 2, 3...)", "Monotonic OCC revision counter. Increments on every PATCH /metadata update. Enforced via DynamoDB ConditionExpressions.", ACCENT_AMBER)
    ]
    cw = Inches(3.75)
    cgap = Inches(0.24)
    cx_start = Inches(0.8)
    cy = Inches(1.65)
    ch = Inches(1.5)

    for i, (name, key, desc, col) in enumerate(counters):
        cx = cx_start + i * (cw + cgap)
        add_card(slide7, cx, cy, cw, ch, CARD_BG, CARD_BORDER)

        tb = slide7.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.12), cw - Inches(0.3), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col

        p_k = tf.add_paragraph()
        p_k.text = key
        p_k.font.name = "Consolas"
        p_k.font.size = Pt(9)
        p_k.font.color.rgb = TEXT_WHITE

        db = slide7.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.65), cw - Inches(0.3), Inches(0.75))
        tf_d = db.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(8.5)
        p_d.font.color.rgb = TEXT_MUTED

    mat_y = Inches(3.35)
    mat_rows = [
        ("Field Names", "Category", "Mutability Rule", "Update & Governance Mechanism"),
        ("document_id, document_class, created_at, created_by", "System Identity", "Strictly Immutable", "Assigned once on initial document creation. Reject modification with 400 error."),
        ("application_version, content_type, content_length, content_checksum", "Binary Specifier", "Immutable for Version", "Cryptographically bound to raw bytes (SHA-256). Updates only on POST /versions."),
        ("metadata_revision, metadata_updated_at, metadata_updated_by", "Concurrency Tracking", "System Monotonic", "Automatically managed by Lambda and DynamoDB OCC conditional expressions."),
        ("customer_id, complete_customer_id_code, account_id, business_area_code, loan_number, currency...", "Domain & Shared Traits", "Mutable via PATCH", "Modifiable via PATCH /metadata with Ajv JSON Schema validation and expected_metadata_revision.")
    ]

    mat_table_shape = slide7.shapes.add_table(len(mat_rows), 4, Inches(0.8), mat_y, Inches(11.733), Inches(3.45))
    mat_table = mat_table_shape.table
    mat_table.columns[0].width = Inches(3.6)
    mat_table.columns[1].width = Inches(2.0)
    mat_table.columns[2].width = Inches(2.3)
    mat_table.columns[3].width = Inches(3.833)

    for r_idx, row in enumerate(mat_rows):
        for c_idx, val in enumerate(row):
            cell = mat_table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(22, 30, 46)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.1)
            tf.margin_top = tf.margin_bottom = Inches(0.06)

            p = tf.paragraphs[0]
            p.text = val
            p.font.name = FONT_FAMILY
            if r_idx == 0:
                p.font.size = Pt(9.5)
                p.font.bold = True
                p.font.color.rgb = AWS_ORANGE if c_idx == 0 else ACCENT_BLUE
            else:
                p.font.size = Pt(8.5)
                if c_idx == 2:
                    p.font.bold = True
                    if "Strictly" in val:
                        p.font.color.rgb = ACCENT_ROSE
                    elif "Immutable for" in val:
                        p.font.color.rgb = ACCENT_AMBER
                    elif "Monotonic" in val:
                        p.font.color.rgb = ACCENT_PURPLE
                    else:
                        p.font.color.rgb = ACCENT_GREEN
                elif c_idx == 0:
                    p.font.name = "Consolas"
                    p.font.color.rgb = TEXT_WHITE
                else:
                    p.font.color.rgb = TEXT_MUTED

    add_footer(slide7, 7, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 8: Metadata Lifecycle & Optimistic Concurrency Control (OCC)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Optimistic Concurrency Control (OCC) & Metadata Evolution",
               "Preventing Lost Overwrites Under High Concurrency & In-Memory Ajv JSON Schema Validation",
               "Concurrency & Evolution")

    add_card(slide8, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    left_b = slide8.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf = left_b.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "OCC Algorithm (PATCH /v1/documents/{id}/metadata)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    occ_steps = [
        ("1. Client Passes Revision:", "Caller includes expected_metadata_revision: N in the request body along with desired field changes."),
        ("2. Read & Consistency Check:", "Lambda reads DynamoDB DOC#{id} pointer. If current_metadata_revision != expected_metadata_revision, aborts immediately with 409 METADATA_CONFLICT."),
        ("3. Fetch S3 Annotation & Merge:", "Fetches current S3 annotation JSON via GetObjectAnnotation, validates that changes contain no immutable fields, increments revision to N+1, and validates against Ajv."),
        ("4. S3 Annotation Commit:", "S3 PutObjectAnnotation attaches updated JSON to S3 object version."),
        ("5. Atomic DynamoDB Conditional Commit:", "DynamoDB executes conditional update (ConditionExpression: current_metadata_revision = :expRev). If parallel update raced, DynamoDB throws ConditionalCheckFailedException.")
    ]
    for title, desc in occ_steps:
        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        p.space_before = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(1)

    add_card(slide8, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    right_b = slide8.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf_r = right_b.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "Metadata Schema Evolution & GitOps"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    p_sch = tf_r.add_paragraph()
    p_sch.text = "Current GitOps Lifecycle (Build-Time Packaging):"
    p_sch.font.name = FONT_FAMILY
    p_sch.font.size = Pt(10.5)
    p_sch.font.bold = True
    p_sch.font.color.rgb = ACCENT_PURPLE
    p_sch.space_before = Pt(6)

    schema_points = [
        ("• Zero Cold Start Latency: ", "Schemas (schemas/loan_agreement-v1.json) are bundled directly into Lambda artifacts and compiled into memory at startup using Ajv. Zero DB network hops."),
        ("• Closed Envelope Contract: ", "additionalProperties: false prevents property pollution and schema drift across document classes."),
        ("• Strict Pull Request Governance: ", "Schema updates require code review, automated schema linting, and CI/CD validation before deployment."),
        ("• Phase 2 Roadmap (Dynamic Registry): ", "Future phase introduces a dedicated DynamoDB Schema Registry table for runtime administrative schema registration without redeployment.")
    ]
    for label, text in schema_points:
        p = tf_r.add_paragraph()
        p.text = label
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_before = Pt(5)

        run = p.add_run()
        run.text = text
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    add_footer(slide8, 8, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 9: Bimodal Ingestion Architecture
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "Bimodal Ingestion Architecture: Inline vs. Direct S3",
               "Optimizing Latency and Eliminating API Gateway 10 MiB Payload Bottlenecks",
               "Ingestion Deep Dive")

    add_card(slide9, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    left_b = slide9.shapes.add_textbox(Inches(1.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf = left_b.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Pathway 1: Inline Binary Upload (≤ 4 MiB)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    p_sub = tf.add_paragraph()
    p_sub.text = "Single HTTP Request for standard documents, receipts, invoices, and IDs."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9.5)
    p_sub.font.color.rgb = TEXT_MUTED
    p_sub.space_before = Pt(4)

    p_box = tf.add_paragraph()
    p_box.text = (
        'POST /v1/documents\n'
        'Headers:\n'
        '  • Content-Type: application/pdf\n'
        '  • X-Content-SHA256: {hex_hash}\n'
        '  • X-Document-Metadata: {base64_json_envelope}\n'
        'Body: Raw Binary Bytes'
    )
    p_box.font.name = "Consolas"
    p_box.font.size = Pt(9)
    p_box.font.color.rgb = ACCENT_GREEN
    p_box.space_before = Pt(8)

    p_adv = tf.add_paragraph()
    p_adv.text = "Execution Characteristics:"
    p_adv.font.name = FONT_FAMILY
    p_adv.font.size = Pt(10.5)
    p_adv.font.bold = True
    p_adv.font.color.rgb = TEXT_WHITE
    p_adv.space_before = Pt(8)

    p_adv_items = [
        "• Total E2E Latency: ~120ms - 250ms synchronous response.",
        "• Atomic Commit: Binary + S3 Annotation + DynamoDB pointer committed simultaneously.",
        "• Zero Client Coordination: 1 request returns 201 Created with complete metadata."
    ]
    for item in p_adv_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(2)

    add_card(slide9, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), CARD_BG, CARD_BORDER)
    right_b = slide9.shapes.add_textbox(Inches(7.05), Inches(1.9), Inches(5.2), Inches(4.7))
    tf_r = right_b.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "Pathway 2: Direct S3 Presigned Upload (> 4 MiB)"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    p_sub_r = tf_r.add_paragraph()
    p_sub_r.text = "Two-Phase Protocol designed for large loan dossiers, mortgage archives & CAD files."
    p_sub_r.font.name = FONT_FAMILY
    p_sub_r.font.size = Pt(9.5)
    p_sub_r.font.color.rgb = TEXT_MUTED
    p_sub_r.space_before = Pt(4)

    p_box_r = tf_r.add_paragraph()
    p_box_r.text = (
        'Step 1: POST /v1/documents/uploads\n'
        '  → Validates metadata, mints 15-min S3 presigned PUT URL.\n'
        'Step 2: PUT {presigned_url}\n'
        '  → Client streams bytes directly to S3 (line rate).\n'
        'Step 3: POST /v1/uploads/{id}/complete\n'
        '  → Lambda validates S3 HeadObject & activates pointer.'
    )
    p_box_r.font.name = "Consolas"
    p_box_r.font.size = Pt(9)
    p_box_r.font.color.rgb = ACCENT_BLUE
    p_box_r.space_before = Pt(8)

    p_adv_r = tf_r.add_paragraph()
    p_adv_r.text = "Execution Characteristics:"
    p_adv_r.font.name = FONT_FAMILY
    p_adv_r.font.size = Pt(10.5)
    p_adv_r.font.bold = True
    p_adv_r.font.color.rgb = TEXT_WHITE
    p_adv_r.space_before = Pt(8)

    p_adv_items_r = [
        "• Bypasses API Gateway 10 MiB Ceiling: Scales to multi-GB uploads.",
        "• Zero Lambda Memory Overhead: Compute handles only metadata coordination.",
        "• Automated Session TTL: Uncompleted sessions cleaned up after 24h via DynamoDB TTL."
    ]
    for item in p_adv_items_r:
        p = tf_r.add_paragraph()
        p.text = item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(2)

    add_footer(slide9, 9, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 10: Exposed REST APIs Comprehensive Catalog
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "Exposed REST APIs: Complete Operations Catalog",
               "15 Production Endpoints Partitioned Across Ingestion, Lifecycle, Metadata, Download & Search",
               "API Surface & Contracts")

    api_catalog = [
        ("Method & Path", "Operation Name", "Target Components", "RBAC Role Required", "Mode"),
        ("GET /health", "Health Check", "DynamoDB, S3, OpenSearch", "Public / All Roles", "Sync"),
        ("POST /v1/documents", "Inline Document Upload", "API GW, Lambda, S3, DynamoDB", "Document.Writer / Admin", "Sync"),
        ("POST /v1/documents/uploads", "Direct Upload Init", "Lambda, S3 Presign, DynamoDB", "Document.Writer / Admin", "Sync"),
        ("POST /v1/uploads/{id}/complete", "Complete Direct Upload", "Lambda, S3 Annotate, DynamoDB", "Document.Writer / Admin", "Sync"),
        ("DELETE /v1/uploads/{id}", "Cancel Direct Upload", "Lambda, S3 Abort, DynamoDB", "Document.Writer / Admin", "Sync"),
        ("GET /v1/documents/{id}", "Get Document Pointer", "API GW, Lambda, DynamoDB", "Document.Reader / All", "Sync"),
        ("GET /v1/documents/{id}/versions", "List Version History", "Lambda, DynamoDB VER# Query", "Document.Reader / All", "Sync"),
        ("POST /v1/documents/{id}/versions", "Create New Version", "Lambda, S3 Binary+Annot, DynamoDB", "Document.Writer / Admin", "Sync"),
        ("GET /v1/documents/{id}/versions/{v}", "Get Historical Version", "Lambda, DynamoDB VER# Item", "Document.Reader / All", "Sync"),
        ("GET /v1/documents/{id}/metadata", "Get Authoritative Metadata", "Lambda, S3 Annotation Get", "Document.Reader / All", "Sync"),
        ("PATCH /v1/documents/{id}/metadata", "OCC Metadata Update", "Lambda, S3 Annotate, DynamoDB", "Document.MetadataEditor / Writer", "Sync"),
        ("GET /v1/documents/{id}/download", "Get Presigned Download URL", "Lambda, S3 Presign (15 min)", "Document.Reader / All", "Sync"),
        ("POST /v1/documents/{id}/soft-delete", "Soft Delete Document", "Lambda, DynamoDB, CDC De-index", "Document.Admin (Exclusive)", "Sync"),
        ("POST /v1/documents/{id}/restore", "Restore Document", "Lambda, DynamoDB, CDC Re-index", "Document.Admin (Exclusive)", "Sync"),
        ("POST /v1/search", "Search Documents", "Lambda, OpenSearch Serverless", "Document.Reader / All", "Sync")
    ]

    api_table_shape = slide10.shapes.add_table(len(api_catalog), 5, Inches(0.8), Inches(1.65), Inches(11.733), Inches(5.1))
    api_table = api_table_shape.table
    api_table.columns[0].width = Inches(2.6)
    api_table.columns[1].width = Inches(2.3)
    api_table.columns[2].width = Inches(3.1)
    api_table.columns[3].width = Inches(2.7)
    api_table.columns[4].width = Inches(1.033)

    for r_idx, row in enumerate(api_catalog):
        for c_idx, val in enumerate(row):
            cell = api_table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(22, 30, 46)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.04)

            p = tf.paragraphs[0]
            p.text = val
            p.font.name = FONT_FAMILY
            if r_idx == 0:
                p.font.size = Pt(9.5)
                p.font.bold = True
                p.font.color.rgb = AWS_ORANGE if c_idx == 0 else ACCENT_BLUE
            else:
                p.font.size = Pt(8.5)
                if c_idx == 0:
                    p.font.name = "Consolas"
                    if "POST" in val:
                        p.font.color.rgb = ACCENT_GREEN
                    elif "GET" in val:
                        p.font.color.rgb = ACCENT_BLUE
                    elif "PATCH" in val:
                        p.font.color.rgb = ACCENT_AMBER
                    elif "DELETE" in val:
                        p.font.color.rgb = ACCENT_ROSE
                    else:
                        p.font.color.rgb = TEXT_WHITE
                elif c_idx == 3:
                    if "Admin" in val:
                        p.font.color.rgb = ACCENT_PURPLE
                    else:
                        p.font.color.rgb = TEXT_WHITE
                else:
                    p.font.color.rgb = TEXT_MUTED

    add_footer(slide10, 10, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 11: Sequence Deep Dive - Inline Binary Upload (POST /v1/documents)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "Sequence Deep Dive: Inline Binary Ingestion (≤ 4 MiB)",
               "API Sequence 02: POST /v1/documents — Atomic Single-Transaction Upload Flow",
               "API Execution Sequence")

    img_02 = "diagrams/api_sequence_diagrams/02_post_documents_inline.jpg"
    if os.path.exists(img_02):
        slide11.shapes.add_picture(img_02, Inches(0.8), Inches(1.65), Inches(7.8), Inches(5.15))

    add_card(slide11, Inches(8.8), Inches(1.65), Inches(3.733), Inches(5.15), CARD_BG, CARD_BORDER)
    sb11 = slide11.shapes.add_textbox(Inches(9.0), Inches(1.85), Inches(3.333), Inches(4.75))
    tf = sb11.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Sequence Breakdown"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    points11 = [
        ("1. Ingress & Auth:", "API Gateway verifies JWT bearer token against Cognito authorizer (requires Document.Writer or Admin)."),
        ("2. Validation & Checksum:", "Lambda calculates SHA-256 over raw binary payload and validates JSON envelope against bank.document-metadata/1 via Ajv in memory."),
        ("3. Storage Commit:", "Uploads binary to S3: documents/{class}/{doc_id} and attaches native S3 Annotation (document-metadata)."),
        ("4. Atomic DynamoDB Commit:", "Executes TransactWriteItems creating pointer DOC#{id}, initial version VER#0000000001, and idempotency lock."),
        ("5. Downstream CDC:", "DynamoDB Streams triggers StreamProcessor Lambda asynchronously for S3 audit logging and SQS indexing.")
    ]
    for title, desc in points11:
        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_before = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(1)

    add_footer(slide11, 11, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 12: Sequence Deep Dive - Direct Upload Initiation & Complete
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_header(slide12, "Sequence Deep Dive: Direct S3 Presigned Upload Lifecycle",
               "API Sequences 03 & 04: POST /v1/documents/uploads & POST /v1/uploads/{id}/complete",
               "API Execution Sequence")

    img_03 = "diagrams/api_sequence_diagrams/03_post_documents_uploads_init.jpg"
    img_04 = "diagrams/api_sequence_diagrams/04_post_uploads_complete.jpg"

    if os.path.exists(img_03):
        slide12.shapes.add_picture(img_03, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_04):
        slide12.shapes.add_picture(img_04, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide12, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b12 = slide12.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf12 = b12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "Two-Phase Upload Protocol Highlights:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    p_sub = tf12.add_paragraph()
    p_sub.text = "• Phase 1 (Init - 03): Validates metadata schema, creates UPLOAD#{id} session in DynamoDB with 24h TTL, and mints 15-min S3 presigned PUT URL.\n• Direct Transfer: Client streams multi-megabyte/gigabyte binary directly to S3 without passing through API Gateway or Lambda.\n• Phase 2 (Complete - 04): Lambda verifies S3 HeadObject (size/checksum), writes S3 Annotation, updates session state to ACTIVE, and commits DynamoDB document pointer."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide12, 12, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 13: Sequence Deep Dive - Upload Abort & Document Retrieval
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13)
    add_header(slide13, "Sequence Deep Dive: Upload Cancellation & Document Inspection",
               "API Sequences 05 & 06: DELETE /v1/uploads/{id} & GET /v1/documents/{id}",
               "API Execution Sequence")

    img_05 = "diagrams/api_sequence_diagrams/05_delete_uploads_cancel.jpg"
    img_06 = "diagrams/api_sequence_diagrams/06_get_documents_document_id.jpg"

    if os.path.exists(img_05):
        slide13.shapes.add_picture(img_05, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_06):
        slide13.shapes.add_picture(img_06, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide13, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b13 = slide13.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf13 = b13.text_frame
    tf13.word_wrap = True

    p = tf13.paragraphs[0]
    p.text = "Cancellation & Inspection Highlights:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_ROSE

    p_sub = tf13.add_paragraph()
    p_sub.text = "• Cancel Upload (05): Aborts direct upload session in DynamoDB (state = ABORTED) and purges uncommitted staging objects from S3 to prevent orphan storage accumulation.\n• Get Document Pointer (06): Executes fast O(1) strongly consistent read on DynamoDB (pk = DOC#{id}, sk = DOC). Returns active version number, S3 version ID, status, and metadata revision in <10ms."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide13, 13, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 14: Sequence Deep Dive - Version Lineage & Binary Updates
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14)
    add_header(slide14, "Sequence Deep Dive: Version Lineage & Binary Updates",
               "API Sequences 07 & 08: GET /v1/documents/{id}/versions & POST /v1/documents/{id}/versions",
               "API Execution Sequence")

    img_07 = "diagrams/api_sequence_diagrams/07_get_documents_versions.jpg"
    img_08 = "diagrams/api_sequence_diagrams/08_post_documents_versions.jpg"

    if os.path.exists(img_07):
        slide14.shapes.add_picture(img_07, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_08):
        slide14.shapes.add_picture(img_08, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide14, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b14 = slide14.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf14 = b14.text_frame
    tf14.word_wrap = True

    p = tf14.paragraphs[0]
    p.text = "Version Lineage & Binary Evolution Mechanics:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE

    p_sub = tf14.add_paragraph()
    p_sub.text = "• Version Query (07): Fast DynamoDB query (pk = DOC#{id} AND begins_with(sk, 'VER#')) returns ordered version history without paginating S3.\n• New Version Upload (08): Uploads replacement binary to same S3 key, creating a new S3 VersionId. Increments application_version (e.g. 1 → 2), resets metadata_revision to 1, recomputes SHA-256 hash, and inserts new VER# item in DynamoDB."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide14, 14, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 15: Sequence Deep Dive - Historical Version Details & S3 Metadata
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15)
    add_header(slide15, "Sequence Deep Dive: Historical Versions & Authoritative Metadata",
               "API Sequences 09 & 10: GET /v1/documents/{id}/versions/{v} & GET /v1/documents/{id}/metadata",
               "API Execution Sequence")

    img_09 = "diagrams/api_sequence_diagrams/09_get_documents_version_id.jpg"
    img_10 = "diagrams/api_sequence_diagrams/10_get_documents_metadata.jpg"

    if os.path.exists(img_09):
        slide15.shapes.add_picture(img_09, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_10):
        slide15.shapes.add_picture(img_10, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide15, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b15 = slide15.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf15 = b15.text_frame
    tf15.word_wrap = True

    p = tf15.paragraphs[0]
    p.text = "Historical & Authoritative Metadata Retrieval:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    p_sub = tf15.add_paragraph()
    p_sub.text = "• Historical Version Details (09): Fetches historical version record (pk = DOC#{id}, sk = VER#{padded_version}) from DynamoDB, returning immutable content length, checksum, and S3 version ID.\n• Authoritative Metadata JSON (10): Resolves active version pointer from DynamoDB and fetches the complete authoritative S3 JSON annotation (document-metadata) directly from S3 via GetObjectAnnotation."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide15, 15, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 16: Sequence Deep Dive - Optimistic Metadata Patch (API 11)
    # =========================================================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16)
    add_header(slide16, "Sequence Deep Dive: Optimistic Concurrency Metadata Patch",
               "API Sequence 11: PATCH /v1/documents/{id}/metadata — Race-Free Metadata Evolution",
               "API Execution Sequence")

    img_11 = "diagrams/api_sequence_diagrams/11_patch_documents_metadata.jpg"
    if os.path.exists(img_11):
        slide16.shapes.add_picture(img_11, Inches(0.8), Inches(1.65), Inches(7.8), Inches(5.15))

    add_card(slide16, Inches(8.8), Inches(1.65), Inches(3.733), Inches(5.15), CARD_BG, CARD_BORDER)
    sb16 = slide16.shapes.add_textbox(Inches(9.0), Inches(1.85), Inches(3.333), Inches(4.75))
    tf = sb16.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Execution Sequence Steps"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER

    points16 = [
        ("1. Read & OCC Check:", "Strongly consistent read on DynamoDB DOC#{id}. Compares expected_metadata_revision with current_metadata_revision (returns 409 Conflict if mismatch)."),
        ("2. S3 Annotation Retrieval:", "Fetches existing JSON annotation from S3 version via GetObjectAnnotation."),
        ("3. Validation & Schema Check:", "Rejects any attempt to modify immutable fields. Increments revision to N+1 and validates complete merged payload via Ajv."),
        ("4. S3 Annotation Update:", "Writes updated JSON annotation in-place via PutObjectAnnotationCommand."),
        ("5. Atomic Conditional Update:", "Updates DynamoDB pointer with ConditionExpression: current_metadata_revision = :expRev, guaranteeing atomicity.")
    ]
    for title, desc in points16:
        p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_before = Pt(5)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(1)

    add_footer(slide16, 16, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 17: Sequence Deep Dive - Download, Soft Delete & Restore
    # =========================================================================
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17)
    add_header(slide17, "Sequence Deep Dive: Secure Download, Soft Delete & Restore",
               "API Sequences 12, 13 & 14: Direct Binary Access & Non-Destructive Lifecycle Governance",
               "API Execution Sequence")

    img_12 = "diagrams/api_sequence_diagrams/12_get_documents_download.jpg"
    img_13 = "diagrams/api_sequence_diagrams/13_post_documents_soft_delete.jpg"

    if os.path.exists(img_12):
        slide17.shapes.add_picture(img_12, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_13):
        slide17.shapes.add_picture(img_13, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide17, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b17 = slide17.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf17 = b17.text_frame
    tf17.word_wrap = True

    p = tf17.paragraphs[0]
    p.text = "Download & Compliance Lifecycle Highlights:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    p_sub = tf17.add_paragraph()
    p_sub.text = "• Secure Direct Download (12): Generates short-lived (15-min) S3 presigned GET URL. Client streams binary directly from S3, eliminating Lambda egress costs.\n• Logical Soft Delete (13): Sets status: 'SOFT_DELETED' in DynamoDB and triggers CDC stream to remove document from OpenSearch index. S3 WORM binary and annotations remain completely untouched for regulatory compliance or restoration (14)."
    p_sub.font.name = FONT_FAMILY
    p_sub.font.size = Pt(9)
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_before = Pt(2)

    add_footer(slide17, 17, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 18: Sequence Deep Dive - OpenSearch Search & Deep Health Probe
    # =========================================================================
    slide18 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide18)
    add_header(slide18, "Sequence Deep Dive: OpenSearch Search & Deep Health Probe",
               "API Sequences 15 & 01: POST /v1/search & GET /health",
               "API Execution Sequence")

    img_15 = "diagrams/api_sequence_diagrams/15_post_search.jpg"
    img_01 = "diagrams/api_sequence_diagrams/01_get_health.jpg"

    if os.path.exists(img_15):
        slide18.shapes.add_picture(img_15, Inches(0.8), Inches(1.65), Inches(5.7), Inches(3.8))
    if os.path.exists(img_01):
        slide18.shapes.add_picture(img_01, Inches(6.8), Inches(1.65), Inches(5.7), Inches(3.8))

    add_card(slide18, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.2), CARD_BG, CARD_BORDER)
    b18 = slide18.shapes.add_textbox(Inches(1.0), Inches(5.68), Inches(11.333), Inches(1.0))
    tf18 = b18.text_frame
    tf18.word_wrap = True

    p = tf18.paragraphs[0]
    p.text = "Search Querying & Observability Highlights:"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    add_footer(slide18, 18, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 19: UI & Backend Integration Architecture
    # =========================================================================
    slide19 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide19)
    add_header(slide19, "Serverless Web Portal & Client Integration Architecture",
               "CloudFront + S3 SPA, Cognito JWT Authentication & Client-Side SHA256 Checksums",
               "Client Integration Architecture")

    ui_img = "diagrams/ui_backend_integration_architecture.png"
    if os.path.exists(ui_img):
        slide19.shapes.add_picture(ui_img, Inches(0.8), Inches(1.65), Inches(8.0), Inches(5.15))

    add_card(slide19, Inches(9.0), Inches(1.65), Inches(3.533), Inches(5.15), CARD_BG, CARD_BORDER)
    sb19 = slide19.shapes.add_textbox(Inches(9.2), Inches(1.85), Inches(3.133), Inches(4.75))
    tf = sb19.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Integration Highlights"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AWS_ORANGE

    ui_points = [
        ("• 1-Click Persona Simulator: ", "Instant token generation for Document.Admin, Document.Writer, Document.Reader, and Compliance Officer."),
        ("• Automatic Token Ingress: ", "Signs all HTTP requests with Cognito OAuth 2.0 / OpenID Connect JWT tokens."),
        ("• Bimodal File Upload UI: ", "Direct S3 presigned upload with client-side SHA256 hashing and inline binary upload (< 4 MiB)."),
        ("• OpenSearch Discovery: ", "Real-time faceted exploration with dynamic status and attribute filters.")
    ]
    for label, desc in ui_points:
        p = tf.add_paragraph()
        p.text = label
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_before = Pt(6)

        run = p.add_run()
        run.text = desc
        run.font.name = FONT_FAMILY
        run.font.size = Pt(9.5)
        run.font.color.rgb = TEXT_MUTED

    add_footer(slide19, 19, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 20: Cost Model & Financial Sizing
    # =========================================================================
    slide20 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide20)
    left_b = slide20.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(5.6), Inches(4.8))
    tf = left_b.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Cognito RBAC Matrix Enforcement"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    rbac_rows = [
        ("Operation", "Reader", "Writer", "Editor", "Admin"),
        ("POST /v1/documents", "❌", "✅", "❌", "✅"),
        ("GET /v1/documents/{id}", "✅", "✅", "✅", "✅"),
        ("GET /v1/download", "✅", "✅", "✅", "✅"),
        ("POST /versions", "❌", "✅", "❌", "✅"),
        ("PATCH /metadata", "❌", "✅", "✅", "✅"),
        ("POST /soft-delete", "❌", "❌", "❌", "✅"),
        ("POST /restore", "❌", "❌", "❌", "✅"),
        ("POST /search", "✅", "✅", "✅", "✅")
    ]

    rbac_table_shape = slide20.shapes.add_table(len(rbac_rows), 5, Inches(1.0), Inches(2.25), Inches(5.6), Inches(4.3))
    rbac_table = rbac_table_shape.table
    rbac_table.columns[0].width = Inches(2.4)
    rbac_table.columns[1].width = Inches(0.8)
    rbac_table.columns[2].width = Inches(0.8)
    rbac_table.columns[3].width = Inches(0.8)
    rbac_table.columns[4].width = Inches(0.8)

    for r_idx, row in enumerate(rbac_rows):
        for c_idx, val in enumerate(row):
            cell = rbac_table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(22, 30, 46)

            tf_c = cell.text_frame
            tf_c.word_wrap = True
            tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = Inches(0.02)
            p = tf_c.paragraphs[0]
            p.text = val
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            p.font.name = FONT_FAMILY
            if r_idx == 0:
                p.font.size = Pt(8.5)
                p.font.bold = True
                p.font.color.rgb = AWS_ORANGE if c_idx == 0 else ACCENT_BLUE
            else:
                p.font.size = Pt(8)
                if c_idx == 0:
                    p.font.name = "Consolas"
                    p.font.color.rgb = TEXT_WHITE
                else:
                    p.font.color.rgb = ACCENT_GREEN if "✅" in val else ACCENT_ROSE

    add_card(slide20, Inches(7.0), Inches(1.7), Inches(5.533), Inches(5.1), CARD_BG, CARD_BORDER)
    right_b = slide20.shapes.add_textbox(Inches(7.25), Inches(1.9), Inches(5.033), Inches(4.7))
    tf_r = right_b.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "Enterprise Cryptography & WORM Controls"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    sec_items = [
        ("AWS KMS Customer Managed Key (CMK):", "Central key (alias/doc-platform-mvp) with automatic annual rotation encrypting S3 buckets, DynamoDB, SQS queues, and CloudWatch log groups."),
        ("WORM Immutability Policy:", "Explicit s3:DeleteObjectVersion DENY in IAM execution policies. Raw object versions cannot be overwritten or destroyed."),
        ("Zero Public S3 Access:", "All S3 buckets enforce S3 Block Public Access with TLS 1.2+ SSL enforcement in bucket policies."),
        ("Micro-Perimeter IAM Roles:", "Each Lambda has a dedicated execution role granting strictly required actions (e.g. Search Lambda cannot write to S3 or DynamoDB).")
    ]
    for title, desc in sec_items:
        p = tf_r.add_paragraph()
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_before = Pt(6)

        p2 = tf_r.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(1)

    add_footer(slide20, 20, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 21: Infrastructure as Code (AWS CDK) & Stack Architecture
    # =========================================================================
    slide21 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide21)
    add_header(slide21, "Infrastructure as Code: AWS CDK v2 Modular Stacks",
               "Deterministic TypeScript CloudFormation Topologies with Zero Circular Dependencies",
               "Deployment Architecture")

    stacks = [
        ("1. SecurityStack", "Cognito User Pool, App Client & KMS CMK Key", "Provisions identity directory, RBAC groups & central encryption key.", ACCENT_BLUE),
        ("2. StorageStack", "S3 Primary Document Bucket & S3 Audit Bucket", "Configures versioning, encryption, CORS & WORM protection policies.", ACCENT_GREEN),
        ("3. SearchStack", "OpenSearch Serverless Collection (documents-v1)", "Sets up vector/search collection, encryption & network security policies.", ACCENT_PURPLE),
        ("4. MessagingStack", "Amazon SQS Indexing Queue & Dead-Letter Queue", "Buffers CDC events with 3x retry limit and CloudWatch alarm triggers.", AWS_ORANGE),
        ("5. ControlPlaneStack", "DynamoDB Table (doc-platform-mvp-control)", "Provisions single-table design with PITR and DynamoDB Streams enabled.", ACCENT_AMBER),
        ("6. ComputeStack", "8 Dedicated AWS Lambda Functions", "Deploys Command, Query, Search, Stream & Indexer microservices.", ACCENT_BLUE),
        ("7. ApiStack", "Amazon API Gateway REST API & Cognito Authorizer", "Wires routes, authorizer, throttling limits, and CORS configuration.", ACCENT_GREEN),
        ("8. ObservabilityStack", "CloudWatch Alarms, Dashboards & Global Tagging", "Provisions DLQ/5xx alarms, operational dashboards, and global resource tags.", ACCENT_ROSE)
    ]

    st_w = Inches(5.7)
    st_h = Inches(1.15)
    st_gap_x = Inches(0.333)
    st_gap_y = Inches(0.12)
    st_x1 = Inches(0.8)
    st_x2 = Inches(6.833)
    st_y_start = Inches(1.7)

    for idx, (name, res, desc, col) in enumerate(stacks):
        col_idx = idx % 2
        row_idx = idx // 2
        sx = st_x1 if col_idx == 0 else st_x2
        sy = st_y_start + row_idx * (st_h + st_gap_y)

        add_card(slide21, sx, sy, st_w, st_h, CARD_BG, CARD_BORDER)

        tb = slide21.shapes.add_textbox(sx + Inches(0.15), sy + Inches(0.1), st_w - Inches(0.3), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_FAMILY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col

        p_r = tf.add_paragraph()
        p_r.text = f"• Resources: {res}"
        p_r.font.name = FONT_FAMILY
        p_r.font.size = Pt(8.5)
        p_r.font.color.rgb = TEXT_WHITE

        p_d = tf.add_paragraph()
        p_d.text = f"• Purpose: {desc}"
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(8)
        p_d.font.color.rgb = TEXT_MUTED

    add_footer(slide21, 21, TOTAL_SLIDES)

    # =========================================================================
    # SLIDE 22: Strategic Architecture Summary & Architectural Decision Records
    # =========================================================================
    slide22 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide22)
    add_header(slide22, "Strategic Architecture Summary & Decision Records",
               "Why This Cloud-Native Blueprint Delivers Unrivaled Durability, Cost-Efficiency & Compliance",
               "Executive Summary & Value")

    pillars = [
        ("1. Zero Idle TCO", "Operational costs scale strictly with business transactions. Dev/Test and staging environments cost $0 when idle, eliminating multi-million dollar annual ECM licenses.", ACCENT_BLUE),
        ("2. Infinite WORM Durability", "S3 11 9s durability combined with immutable S3 Versioning and explicit IAM Deny policies provides airtight compliance against ransomware or accidental deletion.", ACCENT_GREEN),
        ("3. Blast Radius Isolation", "Decoupled CQRS storage architecture ensures search cluster re-indexing or outages never disrupt core document ingestion, metadata updates, or downloads.", ACCENT_PURPLE)
    ]
    pw = Inches(3.75)
    pgap = Inches(0.24)
    px_start = Inches(0.8)
    py = Inches(1.65)
    ph = Inches(1.7)

    for i, (title, desc, col) in enumerate(pillars):
        px = px_start + i * (pw + pgap)
        add_card(slide22, px, py, pw, ph, CARD_BG, CARD_BORDER)

        tb = slide22.shapes.add_textbox(px + Inches(0.15), py + Inches(0.15), pw - Inches(0.3), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col

        db = slide22.shapes.add_textbox(px + Inches(0.15), py + Inches(0.55), pw - Inches(0.3), Inches(1.05))
        tf_d = db.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        p_d.text = desc
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_MUTED

    adr_y = Inches(3.55)
    adr_rows = [
        ("Architectural Decision", "Chosen Approach", "Alternative Rejected", "Core Strategic Justification"),
        ("Metadata Storage Model", "S3 Native Annotations", "S3 Headers (2 KB) / Sidecar Files", "Overcomes 2 KB limit, avoids sidecar file proliferation, ready for S3 Iceberg SQL analytics."),
        ("Transactional Control Plane", "DynamoDB Single-Table", "Aurora PostgreSQL / DocumentDB", "Single-digit millisecond latency, atomic OCC conditional updates, zero server maintenance."),
        ("Search Engine Topology", "OpenSearch Serverless", "Elasticsearch on EC2 / OpenSearch Cluster", "Zero cluster sizing, auto-scaling OCUs, IAM SigV4 authentication, full blast radius isolation."),
        ("Ingestion Protocol", "Bimodal (Inline ≤4MB, Direct >4MB)", "API Gateway Streaming Only", "Bypasses 10 MiB API Gateway ceiling and Lambda memory costs for multi-gigabyte files."),
        ("Metadata Schema Validation", "Build-Time GitOps (Ajv In-Memory)", "Database Schema Lookups on Cold Start", "Zero cold start latency overhead (<1ms validation), strict PR reviews and Git audit trail.")
    ]

    adr_table_shape = slide22.shapes.add_table(len(adr_rows), 4, Inches(0.8), adr_y, Inches(11.733), Inches(3.25))
    adr_table = adr_table_shape.table
    adr_table.columns[0].width = Inches(2.5)
    adr_table.columns[1].width = Inches(2.8)
    adr_table.columns[2].width = Inches(2.8)
    adr_table.columns[3].width = Inches(3.633)

    for r_idx, row in enumerate(adr_rows):
        for c_idx, val in enumerate(row):
            cell = adr_table.cell(r_idx, c_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(22, 30, 46)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.margin_top = tf.margin_bottom = Inches(0.04)

            p = tf.paragraphs[0]
            p.text = val
            p.font.name = FONT_FAMILY
            if r_idx == 0:
                p.font.size = Pt(9.5)
                p.font.bold = True
                p.font.color.rgb = AWS_ORANGE if c_idx == 0 else ACCENT_BLUE
            else:
                p.font.size = Pt(8)
                if c_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = TEXT_WHITE
                elif c_idx == 1:
                    p.font.color.rgb = ACCENT_GREEN
                elif c_idx == 2:
                    p.font.color.rgb = ACCENT_ROSE
                else:
                    p.font.color.rgb = TEXT_MUTED

    add_footer(slide22, 22, TOTAL_SLIDES)

    # -------------------------------------------------------------------------
    # Save Presentation
    # -------------------------------------------------------------------------
    prs.save(output_path)
    print(f"[SUCCESS] Presentation generated successfully: {output_path}")

if __name__ == "__main__":
    build_presentation()
